"""Генерує комплект для захисту дипломної роботи:

* ``../presentation/diploma_presentation.pptx`` – слайди PowerPoint
  з вбудованими діаграмами і speaker notes (підказки доповідачу).
* ``../presentation/speech_script.docx`` – повний текст доповіді,
  розбитий за слайдами (для друку і підготовки до виступу).

Тривалість доповіді – ~7 хвилин, 16 слайдів.
"""

from __future__ import annotations

import sys
from pathlib import Path
from typing import Optional

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Cm as DocxCm, Pt as DocxPt, RGBColor as DocxRGB

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
PRESENTATION_DIR = REPO_ROOT.parent / "presentation"
FIGURES_DIR = REPO_ROOT.parent / "docs" / "figures"
PPTX_PATH = PRESENTATION_DIR / "diploma_presentation.pptx"
DOCX_PATH = PRESENTATION_DIR / "speech_script.docx"


# ============================================================================
# контент усіх слайдів

SLIDES = [
    {
        "title": "Порівняння методів класифікації\nз розширеним алгоритмом CMA-ES",
        "subtitle": ("Випускна кваліфікаційна (бакалаврська) робота\n"
                     "за спеціальністю 122 «Комп'ютерні науки»"),
        "layout": "title",
        "speech": (
            "Шановний голово екзаменаційної комісії, шановні члени комісії. "
            "Представляю кваліфікаційну роботу за темою «Порівняння методів "
            "класифікації з розширеним алгоритмом CMA-ES». Робота поєднує "
            "теоретичне дослідження двох напрямів машинного навчання – "
            "узагальнених адитивних моделей та еволюційної стратегії CMA-ES, "
            "і їхню програмну реалізацію у власному модулі classification_cma_es."
        ),
    },
    {
        "title": "Актуальність і постановка задачі",
        "layout": "bullets",
        "bullets": [
            "–  Задача класифікації табличних даних залишається однією з "
            "найпоширеніших у машинному навчанні",
            "–  Жодний метод не переважає інші на даних довільної природи "
            "(принцип «немає безплатних обідів»)",
            "–  Потрібне емпіричне порівняння методів на сучасних відкритих "
            "датасетах різної природи",
            "–  Окремий напрям – безградієнтні методи оптимізації для "
            "навчання моделей і автоматичного тюнінгу гіперпараметрів",
            "–  Розширення CMA-ES сумішами розподілів обіцяє переваги на "
            "мультимодальних задачах, але потребує емпіричного підтвердження",
        ],
        "speech": (
            "Задача класифікації табличних даних є однією з найбільш "
            "затребуваних у прикладному машинному навчанні. Однак не існує "
            "методу, який би стабільно переважав інші у задачах довільної "
            "природи – це принцип «немає безплатних обідів». Тому потрібне "
            "емпіричне порівняння методів. Окремий інтерес становлять "
            "безградієнтні методи оптимізації, зокрема алгоритм адаптації "
            "коваріаційної матриці CMA-ES, та його розширення сумішами "
            "нормальних розподілів, яке теоретично має переваги на "
            "мультимодальних задачах."
        ),
    },
    {
        "title": "Мета та завдання роботи",
        "layout": "bullets",
        "bullets": [
            "–  Мета: реалізувати програмний модуль для відтворюваного "
            "порівняння методів класифікації та експериментально "
            "дослідити поведінку класичного й розширеного CMA-ES",
            "–  Опанувати теоретичні основи GAM та CMA-ES зі сумішами",
            "–  Реалізувати єдиний інтерфейс для 12 класифікаторів",
            "–  Самостійно реалізувати GAM-класифікатор і розширений CMA-ES",
            "–  Сформувати методику оцінки: Stratified K-Fold, F1, AUC",
            "–  Провести експерименти на трьох сучасних відкритих датасетах",
        ],
        "speech": (
            "Мета роботи – реалізувати програмний модуль, що забезпечує "
            "відтворюване порівняння дванадцяти моделей класифікації, та з "
            "його допомогою експериментально дослідити поведінку класичного "
            "і розширеного варіантів CMA-ES. Завдання включають опанування "
            "теоретичних основ, самостійну реалізацію GAM-класифікатора та "
            "розширеного CMA-ES зі сумішами, формування методики оцінки і "
            "проведення експериментів на трьох датасетах різного характеру."
        ),
    },
    {
        "title": "Зв'язок з теоретичною частиною",
        "layout": "bullets",
        "bullets": [
            "–  Розділ 1 теорії – Узагальнені адитивні моделі (GAM)",
            "–  Розділ 2 теорії – Алгоритм адаптації коваріаційної матриці",
            "–  Програма реалізує обидва напрями у працюючому коді:",
            "      GAM → модуль classifiers/gam.py",
            "      CMA-ES → cma_es.py + mixture_cma_es.py",
            "–  Модель tuned_gam об'єднує обидва теоретичні розділи у "
            "одну робочу реалізацію",
        ],
        "speech": (
            "Теоретична частина роботи містить два розділи: розділ 1 – "
            "узагальнені адитивні моделі, розділ 2 – алгоритм CMA-ES. "
            "Програмну частину побудовано так, щоб обидва теоретичні "
            "розділи були представлені у працюючому коді. GAM реалізовано "
            "як окремий класифікатор; CMA-ES – у двох варіантах: класичному "
            "і розширеному зі сумішами розподілів. Модель tuned_gam, у якій "
            "CMA-ES оптимізує параметри GAM, безпосередньо об'єднує обидва "
            "теоретичні розділи."
        ),
    },
    {
        "title": "Узагальнені адитивні моделі (GAM)",
        "layout": "bullets",
        "bullets": [
            "–  Сума гладких функцій від кожної ознаки замість лінійної комбінації",
            "–  Кожна функція fᵢ(xᵢ) представлена як сума B-сплайнових базисів",
            "–  Функція зв'язку g(·) розширює апарат на задачі класифікації",
            "–  Параметр згладжування λ контролює компроміс між точністю "
            "й перенавчанням",
            "–  Поєднує гнучкість нелінійних моделей з інтерпретованістю "
            "лінійних",
            "–  Власна реалізація через SplineTransformer + LogisticRegression",
        ],
        "speech": (
            "Узагальнена адитивна модель – це сума гладких функцій від "
            "кожної ознаки. У моїй реалізації кожна ознака проходить через "
            "B-сплайнове базисне перетворення, а коефіцієнти підбирає "
            "логістична регресія. Це математично еквівалентно класичному "
            "GAM з логіт-функцією зв'язку, описаному у підрозділі 1.6 "
            "теоретичної частини. Ключова перевага GAM – поєднання гнучкості "
            "нелінійних моделей з інтерпретованістю, властивою лінійним "
            "моделям."
        ),
    },
    {
        "title": "Алгоритм CMA-ES",
        "layout": "bullets",
        "bullets": [
            "–  Еволюційна стратегія з адаптацією коваріаційної матриці",
            "–  Безградієнтний метод – не потребує похідних цільової функції",
            "–  Цикл: семплування → оцінка → відбір → адаптація параметрів",
            "–  Автоматично підлаштовує форму області пошуку під структуру задачі",
            "–  Самоадаптація масштабу пошуку σ та коваріаційної матриці C",
            "–  Стійкий до локальних мінімумів і мультимодальних ландшафтів",
        ],
        "speech": (
            "Алгоритм адаптації коваріаційної матриці CMA-ES – це еволюційна "
            "стратегія без використання похідних. Один цикл складається з "
            "чотирьох кроків: семплування нових кандидатів з багатовимірного "
            "нормального розподілу, оцінка цільової функції, відбір "
            "найкращих, оновлення параметрів. Ключова особливість – "
            "адаптація коваріаційної матриці й масштабу пошуку, що дозволяє "
            "автоматично підлаштовувати форму області пошуку під структуру "
            "цільової функції."
        ),
    },
    {
        "title": "Розширений CMA-ES зі сумішами розподілів",
        "layout": "bullets",
        "bullets": [
            "–  Унімодальний нормальний розподіл замінено сумішшю k нормальних "
            "компонент",
            "–  Параметри суміші (ws, ms, Cs) оновлюються EM-алгоритмом за "
            "найкращою половиною хромосом",
            "–  Самоадаптивний підбір кількості піків за критеріями "
            "кластерного аналізу",
            "–  Очікувана перевага – на мультимодальних задачах оптимізації",
            "–  Власна з нуля реалізація у модулі mixture_cma_es (≈250 рядків)",
            "–  Технічний фікс: момент-усереднення коваріацій (cov_lr) для "
            "уникнення колапсу",
        ],
        "speech": (
            "Розширений варіант CMA-ES замінює унімодальний розподіл "
            "сумішшю k нормальних компонент. Параметри суміші оновлюються "
            "EM-алгоритмом за найкращою половиною хромосом кожної ітерації. "
            "Я реалізував цей алгоритм самостійно у модулі mixture_cma_es. "
            "У процесі реалізації виявив технічний дефект: чисте EM-оновлення "
            "коваріацій на унімодальних задачах призводить до колапсу "
            "дисперсії. Розв'язав це введенням коефіцієнта момент-усереднення "
            "cov_lr за аналогією з rank-µ оновленням класичного CMA-ES."
        ),
    },
    {
        "title": "Дванадцять моделей класифікації",
        "layout": "bullets",
        "bullets": [
            "5 базових: logreg, svm, knn, mlp, gam",
            "2 з CMA-ES як алгоритмом навчання: cma_classic, cma_mixture",
            "5 з автоматичним підбором гіперпараметрів через CMA-ES:",
            "      tuned_logreg, tuned_svm, tuned_knn, tuned_mlp, tuned_gam",
            "–  Усі моделі реалізують єдиний sklearn-сумісний інтерфейс",
            "–  Модель tuned_gam об'єднує обидва теоретичні розділи",
        ],
        "speech": (
            "У програмі реалізовано дванадцять моделей класифікації, "
            "розподілених на три групи. П'ять базових: логістична регресія, "
            "метод опорних векторів, k-найближчих сусідів, багатошарова "
            "нейронна мережа та GAM. Два варіанти CMA-NN, де CMA-ES "
            "виступає алгоритмом навчання ваг нейронної мережі. П'ять "
            "моделей з префіксом tuned_, де CMA-ES автоматично підбирає "
            "гіперпараметри базової моделі. Усі моделі реалізують єдиний "
            "sklearn-сумісний інтерфейс fit/predict/predict_proba."
        ),
    },
    {
        "title": "Архітектура програмного модуля",
        "layout": "bullets",
        "bullets": [
            "–  Python-пакет classifiers з 13 модулями і чіткою ієрархією залежностей",
            "–  Шар даних: data.py, preprocessing.py",
            "–  Шар моделей: models.py, gam.py, cma_es.py, mixture_cma_es.py, "
            "cma_nn.py, hyperparam_tuning.py",
            "–  Шар оцінки: crossval.py, metrics.py, tracking.py",
            "–  Оркестратор: pipeline.py – єдина точка входу run_pipeline",
            "–  ≈2500 рядків коду, 34 модульні тести на pytest – усі проходять",
        ],
        "speech": (
            "Програма побудована як модульний Python-пакет classifiers із "
            "чіткою ієрархією залежностей. Шар підготовки даних відповідає "
            "за завантаження і препроцесинг. Шар моделей містить усі "
            "дванадцять класифікаторів. Шар оцінки виконує перехресну "
            "валідацію, обчислення метрик і MLflow-трекінг. Оркестратор "
            "pipeline.py є єдиною точкою входу для всіх експериментів. "
            "Загальний обсяг коду – близько двох з половиною тисяч рядків, "
            "повністю покритих 34 модульними тестами pytest."
        ),
    },
    {
        "title": "Методика оцінки",
        "layout": "bullets",
        "bullets": [
            "–  Стратифікований поділ на train/test у пропорції 80/20",
            "–  Stratified K-Fold (k=5) перехресна валідація на тренувальній "
            "частині",
            "–  Три класичні метрики класифікації:",
            "      accuracy – частка вірних передбачень",
            "      F1-score – гармонічне середнє precision і recall",
            "      ROC-AUC – площа під ROC-кривою",
            "–  Для мультикласу: F1 weighted-average, AUC у режимі OvR",
        ],
        "speech": (
            "Якість моделей оцінено за стандартною методикою. Спочатку "
            "стратифікований поділ на тренувальну і тестову вибірки у "
            "пропорції 80 на 20 зі збереженням пропорцій класів. На "
            "тренувальній частині додатково – Stratified K-Fold з п'ятьма "
            "фолдами для оцінки стабільності моделі. Обчислюємо три "
            "метрики: accuracy, F1-score та ROC-AUC. Для мультикласових "
            "задач F1 усереднюється за weighted схемою, AUC обчислюється "
            "у стратегії One-vs-Rest."
        ),
    },
    {
        "title": "Сучасні відкриті датасети",
        "layout": "bullets",
        "bullets": [
            "–  PhiUSIIL Phishing URL Dataset (UCI #967, 2024)",
            "      бінарна класифікація, 54 ознаки URL-адреси",
            "–  Steel Plate Defects (UCI #198)",
            "      мультикласова, 7 типів дефектів металу, 27 ознак",
            "–  Default of Credit Card Clients (UCI #350)",
            "      сильна незбалансованість класів 78/22, прогноз дефолту",
            "–  Усі датасети з відкритого UCI ML Repository",
        ],
        "speech": (
            "Для експериментального порівняння обрано три сучасні відкриті "
            "датасети різного характеру. Перший – PhiUSIIL Phishing URL "
            "Dataset 2024 року: бінарна задача відмежування phishing-сайтів "
            "за 54 ознаками URL. Другий – Steel Plate Defects: мультикласова "
            "задача з семи типів дефектів. Третій – Default of Credit Card "
            "Clients: задача прогнозування дефолту з сильною незбалансованістю "
            "класів 78 на 22 відсотки, де метрика accuracy втрачає "
            "інформативність."
        ),
    },
    {
        "title": "Результати: PhiUSIIL",
        "layout": "bullets",
        "bullets": [
            "–  GAM і tuned_gam – ідеальні Accuracy=1.0000, F1=1.0000, AUC=1.0000",
            "–  Логістична регресія і SVM – F1 ≈ 0.998",
            "–  CMA-NN: F1 ≈ 0.991–0.993",
            "–  Задача в значній мірі лінійно роздільна",
            "–  Усі моделі справляються з нею добре",
            "–  GAM-сплайнове перетворення вловлює залишкові нелінійності",
        ],
        "speech": (
            "На датасеті PhiUSIIL більшість моделей досягли значень метрик, "
            "близьких до одиниці. Безумовним лідером стали GAM і tuned_gam "
            "з ідеальними значеннями всіх трьох метрик – на 600 тестових "
            "прикладах не було жодної помилки класифікації. Це пояснюється "
            "тим, що задача в значній мірі лінійно роздільна, і навіть "
            "найпростіша лінійна модель справляється з нею майже ідеально."
        ),
    },
    {
        "title": "Результати: Steel Plate Defects",
        "layout": "bullets",
        "bullets": [
            "–  Лідер: tuned_svm – F1=0.7597, AUC=0.9281",
            "–  tuned_gam – F1=0.7576, базова логістична регресія – F1≈0.73",
            "–  CMA-NN з 30 ітерацій недостатньо для 7 класів і 27 ознак",
            "–  Найскладніше розрізнити Bumps та Other_Faults",
            "–  Підбір C і γ через CMA-ES дає істотне покращення",
            "–  Для збіжності CMA-NN потрібно ≥100 ітерацій",
        ],
        "speech": (
            "На Steel Plate Defects, де задача стає мультикласовою, "
            "найкращим виявився tuned_svm з F1 0.76. Покращення відносно "
            "базового SVM становить майже одну сотку – істотний приріст "
            "для мультикласу. Базова логістична регресія дала F1 близько "
            "0.73 – задача нелінійна. Безградієнтні методи CMA-NN на цій "
            "задачі вимагають значно більшої кількості ітерацій – простір "
            "ваг для 7 класів і 27 ознак стає завеликим для 30 ітерацій."
        ),
    },
    {
        "title": "Результати: Default of Credit Card Clients",
        "layout": "bullets",
        "bullets": [
            "–  Лідери: tuned_mlp (F1=0.4360), tuned_svm (F1=0.4356), "
            "tuned_gam (F1=0.4322)",
            "–  Базовий GAM – найкращий AUC=0.7311",
            "–  cma_mixture помітно випереджає cma_classic за F1 "
            "(0.3333 проти 0.2000)",
            "–  Це перший емпіричний доказ переваги розширеного CMA-ES",
            "–  Незбалансованість робить цільову функцію мультимодальною",
            "–  Підтверджує теоретичне обґрунтування суміші розподілів",
        ],
        "speech": (
            "На найскладнішому датасеті прогнозування дефолту лідери – "
            "tuned-моделі з F1 близько 0.43. Найцікавіший і "
            "найнесподіваніший результат – cma_mixture істотно "
            "випереджає cma_classic: F1 0.33 проти 0.20. Це перший "
            "випадок у дослідженні, коли розширений CMA-ES емпірично "
            "підтверджує свою теоретичну перевагу на задачі з потенційно "
            "мультимодальним ландшафтом цільової функції, спричиненим "
            "сильною незбалансованістю класів."
        ),
    },
    {
        "title": "Графічний інтерфейс на Streamlit",
        "layout": "bullets",
        "bullets": [
            "–  Веб-дашборд для інтерактивного експериментування",
            "–  Бічна панель: вибір датасета, моделей, фолдів, seed, ітерацій",
            "–  Прогрес-бар і лог виконання у реальному часі",
            "–  Інтерактивна таблиця метрик з progress-барами на адаптивній шкалі",
            "–  Графіки часу навчання та збіжності CMA-ES",
            "–  Експорт CSV; інтеграція з MLflow для трекінгу історії",
        ],
        "speech": (
            "Для зручної демонстрації роботи програми реалізовано "
            "графічний веб-дашборд на основі Streamlit. У боковій панелі "
            "користувач послідовно обирає всі параметри запуску. Після "
            "натискання Запустити з'являється прогрес-бар і лог у реальному "
            "часі. Результати виводяться у вигляді інтерактивної таблиці "
            "з кольоровими progress-барами на адаптивній шкалі, що дозволяє "
            "візуально бачити різницю навіть між дуже близькими значеннями. "
            "Усі експерименти автоматично логуються через MLflow."
        ),
    },
    {
        "title": "Загальні висновки",
        "layout": "bullets",
        "bullets": [
            "–  Мету роботи досягнуто, всі завдання виконано",
            "–  Підтверджено: tuned_* стабільно покращують базові моделі "
            "на всіх датасетах",
            "–  GAM-класифікатор увійшов до групи лідерів на 2 з 3 задач",
            "–  Розширений CMA-ES емпірично виправдав теорію на "
            "незбалансованій задачі",
            "–  Жоден метод не переважає на всіх задачах – підтверджено "
            "принцип NFL",
            "–  Перспективи: подальші дослідження mixture-варіанту на "
            "мультимодальних задачах",
        ],
        "speech": (
            "Підсумовуючи. Поставлену мету досягнуто, всі завдання виконано. "
            "Емпірично підтверджено, що автоматичний підбір гіперпараметрів "
            "через CMA-ES стабільно покращує базові моделі. GAM-класифікатор "
            "виправдав теоретичні очікування й увійшов до групи лідерів на "
            "двох з трьох датасетів. Розширений CMA-ES зі сумішами показав "
            "перевагу над класичним на незбалансованій задачі – перше "
            "емпіричне підтвердження теоретичної гіпотези. Жоден метод не "
            "переважає одночасно на всіх задачах, що повністю узгоджується "
            "з принципом «немає безплатних обідів»."
        ),
    },
    {
        "title": "Дякую за увагу!",
        "subtitle": ("Готовий відповісти на ваші запитання\n\n"
                     "github.com/bohdanvolosheniuk1/classification_cma_es"),
        "layout": "title",
        "speech": (
            "Дякую за увагу. Вихідний код роботи опубліковано у відкритому "
            "репозиторії на GitHub. Готовий відповісти на ваші запитання."
        ),
    },
]


# ============================================================================
# pptx генератор

def _set_pptx_text(text_frame, text, *, size=24, bold=False, color=None, align=None):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _add_bullets(text_frame, bullets, *, size=20):
    text_frame.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def _add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    run = p.add_run()
    run.text = text


def _build_title_slide(prs, slide_data):
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)

    # очистити дефолтні placeholder'и
    for shp in list(slide.shapes):
        if shp.has_text_frame:
            shp.text_frame.text = ""

    # ручний заголовок по центру
    title_box = slide.shapes.add_textbox(Cm(1), Cm(6), Cm(23), Cm(4))
    _set_pptx_text(title_box.text_frame, slide_data["title"],
                   size=44, bold=True, color=RGBColor(0x26, 0x46, 0x53),
                   align=PP_ALIGN.CENTER)

    if "subtitle" in slide_data:
        sub_box = slide.shapes.add_textbox(Cm(1), Cm(11), Cm(23), Cm(3))
        _set_pptx_text(sub_box.text_frame, slide_data["subtitle"],
                       size=24, color=RGBColor(0x44, 0x44, 0x44),
                       align=PP_ALIGN.CENTER)

    # автор знизу
    author_box = slide.shapes.add_textbox(Cm(1), Cm(17), Cm(23), Cm(1.5))
    _set_pptx_text(author_box.text_frame,
                   "Богдан Волошенюк   |   ЧНУ ім. Юрія Федьковича   |   2026",
                   size=14, color=RGBColor(0x77, 0x77, 0x77),
                   align=PP_ALIGN.CENTER)

    _add_notes(slide, slide_data["speech"])


def _build_bullets_slide(prs, slide_data):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    for shp in list(slide.shapes):
        if shp.has_text_frame:
            shp.text_frame.text = ""

    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.7), Cm(23), Cm(1.8))
    _set_pptx_text(title_box.text_frame, slide_data["title"],
                   size=30, bold=True, color=RGBColor(0x26, 0x46, 0x53))

    bullets_box = slide.shapes.add_textbox(Cm(1.5), Cm(3.2), Cm(22), Cm(13))
    _add_bullets(bullets_box.text_frame, slide_data["bullets"], size=20)

    _add_notes(slide, slide_data["speech"])


def _build_image_slide(prs, slide_data):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    for shp in list(slide.shapes):
        if shp.has_text_frame:
            shp.text_frame.text = ""

    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.7), Cm(23), Cm(1.8))
    _set_pptx_text(title_box.text_frame, slide_data["title"],
                   size=30, bold=True, color=RGBColor(0x26, 0x46, 0x53))

    img_path = FIGURES_DIR / slide_data["image"]
    if img_path.exists():
        # вставка з центруванням
        slide_w = prs.slide_width
        target_w = Cm(20)
        left = (slide_w - target_w) // 2
        top = Cm(3.5)
        slide.shapes.add_picture(str(img_path), left, top, width=target_w)
    else:
        msg = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(20), Cm(2))
        _set_pptx_text(msg.text_frame,
                       f"[не знайдено зображення: {slide_data['image']}]",
                       size=18, color=RGBColor(0xaa, 0x00, 0x00))

    _add_notes(slide, slide_data["speech"])


def build_pptx(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Cm(25.4)   # стандарт widescreen 16:9
    prs.slide_height = Cm(19.05)

    for slide_data in SLIDES:
        layout = slide_data.get("layout", "bullets")
        if layout == "title":
            _build_title_slide(prs, slide_data)
        elif layout == "image":
            _build_image_slide(prs, slide_data)
        else:
            _build_bullets_slide(prs, slide_data)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


# ============================================================================
# docx генератор (текст доповіді)

def _docx_set_run(run, *, size=12, bold=False, italic=False,
                  font="Times New Roman", color=None):
    run.font.name = font
    run.font.size = DocxPt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    rfonts.set(qn("w:ascii"), font)
    rfonts.set(qn("w:hAnsi"), font)
    rfonts.set(qn("w:cs"), font)


def _docx_para(doc, text, *, size=12, bold=False, italic=False,
               align=None, indent=DocxCm(1.25), space=DocxPt(6)):
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    pf = p.paragraph_format
    pf.first_line_indent = indent
    pf.space_after = space
    pf.line_spacing = 1.5
    run = p.add_run(text)
    _docx_set_run(run, size=size, bold=bold, italic=italic)
    return p


def _docx_heading(doc, text, level=1):
    sizes = {1: 16, 2: 14}
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT if level > 1 else WD_ALIGN_PARAGRAPH.CENTER
    pf = p.paragraph_format
    pf.space_before = DocxPt(14)
    pf.space_after = DocxPt(8)
    pf.line_spacing = 1.15
    pf.keep_with_next = True
    run = p.add_run(text)
    _docx_set_run(run, size=sizes.get(level, 12), bold=True)


def build_docx(out_path: Path) -> None:
    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = DocxPt(12)
    for section in doc.sections:
        section.left_margin = DocxCm(2.5)
        section.right_margin = DocxCm(1.5)
        section.top_margin = DocxCm(2)
        section.bottom_margin = DocxCm(2)

    # титулка
    for _ in range(3):
        doc.add_paragraph()
    _docx_para(doc, "ТЕКСТ ДОПОВІДІ", size=22, bold=True,
               align=WD_ALIGN_PARAGRAPH.CENTER, indent=None,
               space=DocxPt(12))
    _docx_para(doc, "до захисту дипломної роботи бакалавра",
               size=14, italic=True, align=WD_ALIGN_PARAGRAPH.CENTER,
               indent=None, space=DocxPt(20))
    _docx_para(doc,
        "Порівняння методів класифікації з розширеним алгоритмом CMA-ES",
        size=14, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER,
        indent=None, space=DocxPt(60))
    _docx_para(doc, "Виконав: Волошенюк Богдан Анатолійович, група 441-А",
               size=12, align=WD_ALIGN_PARAGRAPH.CENTER,
               indent=None, space=DocxPt(4))
    _docx_para(doc, "Спеціальність 122 «Комп'ютерні науки»",
               size=12, align=WD_ALIGN_PARAGRAPH.CENTER,
               indent=None, space=DocxPt(4))
    _docx_para(doc, "ЧНУ ім. Юрія Федьковича, 2026",
               size=12, align=WD_ALIGN_PARAGRAPH.CENTER,
               indent=None, space=DocxPt(20))

    doc.add_page_break()

    # інструкція доповідачу
    _docx_heading(doc, "Інструкція доповідачу", level=1)
    _docx_para(doc,
        "Загальна тривалість доповіді – приблизно сім–вісім хвилин. На кожен "
        "слайд відведено від 25 до 45 секунд. Говорити спокійним темпом, "
        "приблизно 130 слів на хвилину. Перед виступом одноразово прочитати "
        "вголос – щоб відчути ритм.")
    _docx_para(doc,
        "Перед слайдом про графічний інтерфейс бажано підготувати "
        "відкритий заздалегідь дашборд у браузері – для живої демонстрації. "
        "Якщо проєктор слабкий або немає інтернету – обмежитись описом за "
        "слайдом.")
    _docx_para(doc,
        "Питання комісії очікувано стосуватимуться: (а) у чому новизна "
        "розширеного варіанту CMA-ES і коли він має перевагу; "
        "(б) чому CMA-NN не завжди перемагає sklearn-моделі; "
        "(в) як обиралися датасети; "
        "(г) як модель tuned_gam об'єднує теоретичні розділи.")

    doc.add_page_break()

    # текст по слайдах
    _docx_heading(doc, "Текст по слайдах", level=1)
    for i, slide_data in enumerate(SLIDES, start=1):
        _docx_heading(doc, f"Слайд {i}. {slide_data['title']}", level=2)

        # підказки що на слайді
        if slide_data.get("layout") == "image":
            _docx_para(doc, f"На слайді: схема ({slide_data['image']})",
                       italic=True, indent=None, space=DocxPt(4))
        elif "bullets" in slide_data:
            _docx_para(doc, "На слайді: маркований список:",
                       italic=True, indent=None, space=DocxPt(2))
            for b in slide_data["bullets"]:
                p = doc.add_paragraph(style="List Bullet")
                p.paragraph_format.space_after = DocxPt(2)
                p.paragraph_format.line_spacing = 1.15
                run = p.add_run(b)
                _docx_set_run(run, size=11)
            doc.add_paragraph()
        elif "subtitle" in slide_data:
            _docx_para(doc, f"На слайді: заголовок і підзаголовок "
                            f"({slide_data['subtitle']})",
                       italic=True, indent=None, space=DocxPt(4))

        # сама репліка
        _docx_para(doc, "Доповідь:", size=12, bold=True,
                   indent=None, space=DocxPt(2))
        _docx_para(doc, slide_data["speech"], size=12)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))


def main() -> int:
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)

    print("Генерую PowerPoint...")
    build_pptx(PPTX_PATH)
    pptx_kb = PPTX_PATH.stat().st_size // 1024
    print(f"  OK: {PPTX_PATH} ({pptx_kb} KB)")

    print("Генерую текст доповіді (.docx)...")
    build_docx(DOCX_PATH)
    docx_kb = DOCX_PATH.stat().st_size // 1024
    print(f"  OK: {DOCX_PATH} ({docx_kb} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())

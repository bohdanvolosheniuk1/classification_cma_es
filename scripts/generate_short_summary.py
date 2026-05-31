"""Коротка презентація (5 слайдів) — що зроблено по проєкту.

Сучасний двоколонковий layout 16:9: ліворуч — буллет-пункти,
праворуч — підходяща діаграма. Кольорова палітра:

* основний:   #264653 — темний бірюзовий
* акцент 1:   #2A9D8F — м'ятний
* акцент 2:   #E76F51 — кораловий
* фон:        білий
* текст:      #222222

Файли:
* ``../../presentation/short_summary.pptx`` — сама презентація
* ``../../presentation/short_summary_speech.docx`` — текст доповіді
"""

from __future__ import annotations

import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH as DOCX_ALIGN
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Cm as DCm, Pt as DPt

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = REPO_ROOT.parent / "presentation"
PPTX_PATH = OUT_DIR / "short_summary.pptx"
DOCX_PATH = OUT_DIR / "short_summary_speech.docx"

DOCS_FIG = REPO_ROOT.parent / "docs" / "figures"
DIPLOMA_FIG = REPO_ROOT.parent / "diploma_figures"


# Кольорова палітра
COLOR_PRIMARY = RGBColor(0x26, 0x46, 0x53)
COLOR_ACCENT_1 = RGBColor(0x2A, 0x9D, 0x8F)
COLOR_ACCENT_2 = RGBColor(0xE7, 0x6F, 0x51)
COLOR_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)


SLIDES = [
    {
        "kind": "title",
        "title": "Дипломна робота бакалавра",
        "subtitle": "classification_cma_es:\n"
                    "порівняння класифікаторів з розширеним CMA-ES",
        "footer": "Богдан Волошенюк   |   ЧНУ ім. Юрія Федьковича   |   2026",
        "speech": (
            "Вітаю. Представляю короткий підсумок виконаної роботи по "
            "дипломному проєкту бакалавра. За цією презентацією — "
            "повний обсяг зробленого: теоретична частина, програмна "
            "реалізація, документація і фінальні матеріали для захисту. "
            "Тема — порівняння методів класифікації з застосуванням "
            "розширеного алгоритму CMA-ES."
        ),
    },
    {
        "kind": "two_col",
        "number": "01",
        "title": "Теоретична частина",
        "bullets": [
            "Розділ 1 — Узагальнені адитивні моделі (GAM)",
            "6 підрозділів: від базових понять до функцій зв'язку",
            "Розділ 2 — Алгоритм CMA-ES і його розширення",
            "Опрацьовано дисертацію Літвінчук Ю.А. (2024)",
            "~30 сторінок основного тексту",
        ],
        "image": DOCS_FIG / "05_cmaes_cycle.png",
        "speech": (
            "Теоретична частина містить два повноцінні розділи. Перший — "
            "узагальнені адитивні моделі GAM, шість підрозділів від "
            "загальних понять до функцій зв'язку. Другий — алгоритм "
            "CMA-ES і його розширення зі сумішами нормальних розподілів, "
            "теж шість підрозділів. У роботі ґрунтовно опрацьовано "
            "матеріал дисертації Літвінчук Юлії Анатоліївни 2024 року. "
            "Загальний обсяг теоретичної частини — близько тридцяти "
            "сторінок з висновками по кожному розділу."
        ),
    },
    {
        "kind": "two_col",
        "number": "02",
        "title": "Програмна реалізація",
        "bullets": [
            "12 моделей: LogReg, SVM, kNN, MLP, GAM",
            "+ CMA-NN classic / mixture (5-й метод за Літвінчук)",
            "+ 5 моделей tuned_* (підбір параметрів через CMA-ES)",
            "3 датасети: PhiUSIIL, Steel Plate, Loan Approval",
            "Streamlit GUI, CLI, MLflow, 34 тести pytest",
            "Відкритий код на GitHub",
        ],
        "image": DOCS_FIG / "01_architecture.png",
        "speech": (
            "На основі теорії розроблено повноцінний програмний засіб "
            "classification_cma_es. Усього в програмі дванадцять моделей "
            "класифікації: п'ять класичних, два варіанти CMA-NN — це і є "
            "п'ятий метод за дисертацією Літвінчук — та п'ять моделей з "
            "автоматичним підбором гіперпараметрів через CMA-ES. "
            "Програма працює з трьома сучасними датасетами. Особлива "
            "увага приділена реалізації розширеного CMA-ES з власною "
            "технічною поправкою cov_lr. Інтерфейси — командний рядок "
            "і графічний дашборд на Streamlit. Експерименти журналюються "
            "через MLflow. Усе покрито тридцятьма чотирма тестами і "
            "викладено у відкритий репозиторій на GitHub."
        ),
    },
    {
        "kind": "two_col",
        "number": "03",
        "title": "Розділи 3-4 диплома + аналіз",
        "bullets": [
            "Розділ 3 — Програмна реалізація (8 підрозділів)",
            "Розділ 4 — Результати експериментів (5 підрозділів)",
            "7 рисунків з реальних прогонів моделей",
            "9 математичних формул (GAM, EM, CMA-ES, метрики)",
            "3 скрини дашборду + зведений графік 12 моделей",
        ],
        "image": DIPLOMA_FIG / "comparison_f1.png",
        "speech": (
            "До диплома додано два нові розділи. Третій — програмна "
            "реалізація, вісім підрозділів. Четвертий — результати "
            "експериментів, п'ять підрозділів. Для ілюстрації виконано "
            "сім рисунків з реальних прогонів моделей: матриці плутанини, "
            "ROC-криві, коваріаційна матриця CMA-ES, криві збіжності. "
            "У текст вставлено дев'ять математичних формул, три "
            "скриншоти дашборду та зведений графік порівняння всіх "
            "дванадцяти моделей на трьох датасетах."
        ),
    },
    {
        "kind": "two_col",
        "number": "04",
        "title": "Готовність до захисту",
        "bullets": [
            "Зміст з номерами сторінок + бібліографія (13 джерел)",
            "Фінальний диплом: 657 KB, 456 параграфів, 17 рисунків",
            "PDF-гайд, HTML API, Word-документація",
            "Дві презентації: повна (16 слайдів) + ця коротка",
            "GitHub: github.com/bohdanvolosheniuk1/classification_cma_es",
        ],
        "image": DIPLOMA_FIG / "dashboard_results.png",
        "speech": (
            "Усе оформлено за стандартами. Зміст із номерами сторінок, "
            "бібліографія на тринадцять джерел у форматі ДСТУ. "
            "Фінальний документ диплома — близько шестисот шістдесяти "
            "кілобайтів, чотириста п'ятдесят шість параграфів і "
            "сімнадцять рисунків. Окремо створено PDF-гайд, HTML "
            "API-документацію через pdoc, формальну Word-документацію "
            "та дві презентації для захисту. Усі матеріали готові, "
            "проєкт можна здавати куратору."
        ),
    },
]


# ============================================================================
# helpers — PPTX

def _shape_rect(slide, x, y, w, h, fill, line_color=None):
    """Прямокутник-плашка для фону блоку."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def _text(slide, x, y, w, h, text, *, size=18, bold=False, color=None,
          align=None, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return tf


def _bullets(slide, x, y, w, h, bullets, *, size=18, color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_top = Cm(0.05)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        # маркер
        m_run = p.add_run()
        m_run.text = "■  "
        m_run.font.name = "Calibri"
        m_run.font.size = Pt(size)
        m_run.font.color.rgb = COLOR_ACCENT_1
        # сам текст
        t_run = p.add_run()
        t_run.text = line
        t_run.font.name = "Calibri"
        t_run.font.size = Pt(size)
        t_run.font.color.rgb = color or COLOR_TEXT


def _add_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    run = p.add_run()
    run.text = text


# ============================================================================
# slides

def _build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # лівий кольоровий бар
    _shape_rect(slide, Cm(0), Cm(0), Cm(1.5), prs.slide_height,
                fill=COLOR_PRIMARY)

    # центральна частина — заголовок
    _text(slide, Cm(2), Cm(4.5), Cm(22.5), Cm(2.5),
          data["title"], size=42, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

    # лінія-розділювач
    _shape_rect(slide, Cm(10), Cm(7.5), Cm(5.4), Cm(0.1),
                fill=COLOR_ACCENT_2)

    # підзаголовок
    _text(slide, Cm(2), Cm(8.2), Cm(22.5), Cm(2.5),
          data["subtitle"], size=22,
          color=COLOR_TEXT, align=PP_ALIGN.CENTER)

    # футер
    _text(slide, Cm(2), Cm(12.8), Cm(22.5), Cm(1.5),
          data["footer"], size=14,
          color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    _add_speaker_notes(slide, data["speech"])


def _build_two_col_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # шапка з номером і назвою
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, Cm(2.3),
                fill=COLOR_PRIMARY)
    _text(slide, Cm(0.7), Cm(0.3), Cm(2), Cm(1.8),
          data["number"], size=32, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)
    _text(slide, Cm(2.8), Cm(0.55), Cm(20), Cm(1.5),
          data["title"], size=26, bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT)

    # ліва колонка — буллети
    _bullets(slide, Cm(0.7), Cm(3), Cm(11.5), Cm(10),
             data["bullets"], size=18)

    # права колонка — картинка
    image_path = data.get("image")
    if image_path and Path(image_path).exists():
        # рамка під картинкою
        _shape_rect(slide, Cm(12.7), Cm(2.8), Cm(12.0), Cm(10.4),
                    fill=COLOR_BG_LIGHT)
        # сама картинка
        slide.shapes.add_picture(str(image_path),
                                 Cm(13.0), Cm(3.1),
                                 width=Cm(11.4), height=Cm(9.8))

    # підвал слайду
    _shape_rect(slide, Cm(0), Cm(13.7), prs.slide_width, Cm(0.6),
                fill=COLOR_ACCENT_1)

    _add_speaker_notes(slide, data["speech"])


def build_pptx() -> Path:
    prs = Presentation()
    # 16:9
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.29)

    for data in SLIDES:
        if data["kind"] == "title":
            _build_title_slide(prs, data)
        else:
            _build_two_col_slide(prs, data)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_PATH))
    return PPTX_PATH


# ============================================================================
# DOCX — текст доповіді

def _set_run(run, *, size=12, bold=False, italic=False,
             font="Times New Roman"):
    run.font.name = font
    run.font.size = DPt(size)
    run.font.bold = bold
    run.font.italic = italic
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    rfonts.set(qn("w:ascii"), font)
    rfonts.set(qn("w:hAnsi"), font)
    rfonts.set(qn("w:cs"), font)


def _docx_para(doc, text, *, size=12, bold=False, italic=False, align=None,
               indent=DCm(1.25), space=DPt(6)):
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    pf = p.paragraph_format
    if indent is not None:
        pf.first_line_indent = indent
    pf.space_after = space
    pf.line_spacing = 1.5
    run = p.add_run(text)
    _set_run(run, size=size, bold=bold, italic=italic)


def _docx_heading(doc, text, level=1):
    p = doc.add_paragraph()
    pf = p.paragraph_format
    pf.space_before = DPt(14)
    pf.space_after = DPt(8)
    pf.keep_with_next = True
    pf.line_spacing = 1.15
    p.alignment = DOCX_ALIGN.CENTER if level == 1 else DOCX_ALIGN.LEFT
    run = p.add_run(text)
    _set_run(run, size=16 if level == 1 else 13, bold=True)


def build_docx() -> Path:
    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = DPt(12)
    for section in doc.sections:
        section.left_margin = DCm(2.5)
        section.right_margin = DCm(1.5)
        section.top_margin = DCm(2)
        section.bottom_margin = DCm(2)

    # титулка
    for _ in range(4):
        doc.add_paragraph()
    _docx_para(doc, "ТЕКСТ ДОПОВІДІ",
               size=20, bold=True, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(8))
    _docx_para(doc, "Коротка презентація (5 слайдів)",
               size=14, italic=True, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(20))
    _docx_para(doc, "classification_cma_es",
               size=14, bold=True, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(60))
    _docx_para(doc, "Богдан Волошенюк",
               size=12, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(2))
    _docx_para(doc, "ЧНУ ім. Юрія Федьковича, 2026",
               size=12, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(20))
    doc.add_page_break()

    # короткі підказки
    _docx_heading(doc, "Для доповідача", level=1)
    _docx_para(doc,
        "Загальна тривалість — приблизно 2.5–3 хвилини. На кожний "
        "слайд відведено від 20 до 40 секунд. Темп мовлення спокійний, "
        "близько 130 слів на хвилину. Перед виступом одноразово "
        "прочитати вголос — щоб відчути ритм. Тримати презентацію у "
        "Режимі доповідача (Presenter View) — там видно нотатки до "
        "кожного слайда.")
    doc.add_page_break()

    _docx_heading(doc, "Текст по слайдах", level=1)
    for i, data in enumerate(SLIDES, start=1):
        _docx_heading(doc, f"Слайд {i}. {data['title']}", level=2)

        if data["kind"] == "two_col":
            _docx_para(doc, "На слайді:", bold=True, indent=None, space=DPt(2))
            for b in data["bullets"]:
                p = doc.add_paragraph()
                pf = p.paragraph_format
                pf.left_indent = DCm(0.75)
                pf.first_line_indent = DCm(-0.5)
                pf.space_after = DPt(3)
                pf.line_spacing = 1.4
                run = p.add_run("•  " + b)
                _set_run(run, size=11)
            img_path = data.get("image")
            if img_path:
                _docx_para(doc, f"Праворуч — зображення: {Path(img_path).name}",
                           italic=True, indent=None, space=DPt(4))
        elif data["kind"] == "title":
            _docx_para(doc, f"На слайді: заголовок «{data['title']}» + "
                            f"підзаголовок.", italic=True,
                       indent=None, space=DPt(4))

        _docx_para(doc, "Доповідь:", bold=True, indent=None, space=DPt(2))
        _docx_para(doc, data["speech"], size=12)
        doc.add_paragraph()

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(DOCX_PATH))
    return DOCX_PATH


def main() -> int:
    print("Генерую PowerPoint...")
    pptx_out = build_pptx()
    print(f"  OK: {pptx_out} ({pptx_out.stat().st_size // 1024} KB)")

    print("Генерую текст доповіді (.docx)...")
    docx_out = build_docx()
    print(f"  OK: {docx_out} ({docx_out.stat().st_size // 1024} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())

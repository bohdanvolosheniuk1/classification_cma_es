"""Презентація для усного виступу (3-5 хв, 6 слайдів) + текст доповіді.

На відміну від двох попередніх презентацій:
* ``diploma_presentation.pptx`` — повна для захисту (16 слайдів, 7 хв).
* ``short_summary.pptx`` — підсумкова показова (5 слайдів).
* **ця**: ``oral_presentation.pptx`` — для усного виступу (6 слайдів,
  ~4 хв розповіді).

Включає окремий текст для зачитування у файлі
``oral_presentation_speech.docx``.
"""

from __future__ import annotations

import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH as DOCX_ALIGN
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Cm as DCm, Pt as DPt

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = REPO_ROOT.parent / "presentation"
PPTX_PATH = OUT_DIR / "oral_presentation.pptx"
DOCX_PATH = OUT_DIR / "oral_presentation_speech.docx"

DOCS_FIG = REPO_ROOT.parent / "docs" / "figures"
DIPLOMA_FIG = REPO_ROOT.parent / "diploma_figures"


# Кольорова палітра (сучасна, спокійна)
COLOR_PRIMARY = RGBColor(0x26, 0x46, 0x53)   # темний бірюзовий
COLOR_ACCENT_1 = RGBColor(0x2A, 0x9D, 0x8F)  # м'ятний
COLOR_ACCENT_2 = RGBColor(0xE7, 0x6F, 0x51)  # кораловий
COLOR_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


SLIDES = [
    # ---- 1. Title ----
    {
        "kind": "title",
        "title": "Порівняння методів класифікації\nз розширеним алгоритмом CMA-ES",
        "subtitle": "Дипломна робота бакалавра",
        "footer": "Богдан Волошенюк   |   ЧНУ ім. Юрія Федьковича   |   2026",
        "speech": (
            "Добрий день, шановні члени комісії. Представляю дипломну "
            "роботу на тему порівняльного аналізу методів класифікації "
            "машинного навчання із застосуванням розширеного алгоритму "
            "CMA-ES. Робота поєднує теоретичне дослідження двох "
            "напрямків машинного навчання та повноцінну програмну "
            "реалізацію з експериментальним порівнянням."
        ),
    },
    # ---- 2. Theory ----
    {
        "kind": "two_images",
        "number": "01",
        "title": "Теоретична частина",
        "bullets": [
            "Розділ 1 — Узагальнені адитивні моделі (GAM): гнучкі "
            "моделі через сплайнові базисні функції",
            "Розділ 2 — Алгоритм CMA-ES: еволюційна оптимізація "
            "без похідних",
            "Розширення зі сумішами розподілів за дисертацією "
            "Літвінчук Ю.А. (2024)",
        ],
        "image_left": DOCS_FIG / "04_gam.png",
        "image_right": DOCS_FIG / "05_cmaes_cycle.png",
        "speech": (
            "Теоретична частина містить два розділи. Перший присвячено "
            "узагальненим адитивним моделям GAM. Це гнучкі моделі, які "
            "поєднують інтерпретованість лінійної регресії з можливістю "
            "моделювати нелінійні залежності через сплайнові базисні "
            "функції. Другий розділ — алгоритм адаптації коваріаційної "
            "матриці CMA-ES. Це сучасний еволюційний оптимізатор, який "
            "не потребує обчислення похідних. У роботі також ґрунтовно "
            "розглянуто його розширення зі сумішами нормальних "
            "розподілів, запропоноване у дисертації Літвінчук Юлії "
            "Анатоліївни у 2024 році."
        ),
    },
    # ---- 3. Data ----
    {
        "kind": "two_col",
        "number": "02",
        "title": "Дані для дослідження",
        "bullets": [
            "PhiUSIIL Phishing URL — UCI #967, 2024 рік",
            "235 тис. URL-адрес, бінарна класифікація",
            "Steel Plate Defects — UCI #198",
            "7 типів дефектів металевих пластин, мультиклас",
            "Default of Credit Card Clients — UCI #350",
            "30 тис. клієнтів, незбалансовані класи 78 / 22",
        ],
        "image": DIPLOMA_FIG / "phiusiil_roc.png",
        "speech": (
            "Для дослідження обрано три набори даних з UCI Machine "
            "Learning Repository. Перший — PhiUSIIL Phishing URL "
            "Dataset 2024 року, бінарна класифікація 235 тисяч URL-"
            "адрес на phishing і легітимні. Другий — Steel Plate "
            "Defects, мультикласова задача з семи типів дефектів "
            "металевих пластин. Третій — Default of Credit Card "
            "Clients, бінарна задача прогнозування дефолту з сильно "
            "незбалансованими класами у пропорції 78 на 22 відсотки. "
            "Така різноманітність дозволяє побачити, як методи "
            "поводяться у принципово різних умовах."
        ),
    },
    # ---- 4. Programming ----
    {
        "kind": "two_col",
        "number": "03",
        "title": "Програмна реалізація",
        "bullets": [
            "Модуль classification_cma_es на Python",
            "12 моделей: 5 базових + 2 CMA-NN + 5 tuned_*",
            "Власна реалізація розширеного CMA-ES за [5]",
            "GAM-класифікатор без сторонніх бібліотек",
            "Streamlit GUI + CLI + MLflow + 34 тести pytest",
            "Відкритий код на GitHub",
        ],
        "image": DOCS_FIG / "01_architecture.png",
        "speech": (
            "На основі теорії розроблено програмний модуль "
            "classification_cma_es мовою Python. Усього в програмі "
            "реалізовано дванадцять моделей класифікації: п'ять "
            "класичних — це логістична регресія, метод опорних "
            "векторів, k-найближчих сусідів, нейронна мережа та власна "
            "реалізація узагальненої адитивної моделі GAM. Два варіанти "
            "CMA-NN — нейронної мережі, вагові параметри якої "
            "оптимізуються алгоритмом CMA-ES, і це той самий п'ятий "
            "метод за дисертацією Літвінчук. Та п'ять моделей з "
            "префіксом tuned_, в яких CMA-ES автоматично підбирає "
            "гіперпараметри. Програма має командний інтерфейс і "
            "графічний дашборд на Streamlit. Покрита тридцятьма "
            "чотирма тестами на pytest. Весь код викладено на GitHub."
        ),
    },
    # ---- 5. Results ----
    {
        "kind": "big_image",
        "number": "04",
        "title": "Результати порівняння 12 моделей",
        "image": DIPLOMA_FIG / "comparison_f1.png",
        "caption": "F1-score 12 моделей на трьох датасетах",
        "footer_bullets": [
            "GAM з ідеальними метриками на PhiUSIIL",
            "tuned_svm лідирує на Steel Plate",
            "tuned_mlp найкращий на Loan Approval",
            "Підбір параметрів через CMA-ES — стабільне покращення",
        ],
        "speech": (
            "Експерименти показали наочний результат. Універсального "
            "лідера серед методів немає — на кожному датасеті виграє "
            "інший метод. На PhiUSIIL найкращим виявився GAM з "
            "ідеальними значеннями всіх метрик. На Steel Plate Defects "
            "лідирує метод опорних векторів з підбором гіперпараметрів. "
            "На незбалансованому Loan Approval — нейронна мережа з "
            "tuned-параметрами. Загальний висновок: автоматичний "
            "підбір гіперпараметрів через CMA-ES стабільно покращує "
            "базові моделі на всіх трьох датасетах. А розширений "
            "CMA-ES зі сумішами показав цікаве покращення F1-score "
            "саме на незбалансованій задачі — що підтверджує теоретичну "
            "гіпотезу з дисертації."
        ),
    },
    # ---- 6. Conclusion ----
    {
        "kind": "two_col",
        "number": "05",
        "title": "Висновок і власний внесок",
        "bullets": [
            "Обидва теоретичні розділи втілено в коді",
            "tuned_gam об'єднує GAM і CMA-ES в одну модель",
            "Власна поправка cov_lr усуває колапс дисперсії в [5]",
            "Повний пакет документації: PDF, HTML API, презентації",
            "GitHub: bohdanvolosheniuk1/classification_cma_es",
            "Дякую за увагу!",
        ],
        "image": DIPLOMA_FIG / "dashboard_results.png",
        "speech": (
            "Підбиваючи підсумки. У роботі повноцінно реалізовано "
            "обидва теоретичні розділи — GAM як окремий класифікатор і "
            "розширений CMA-ES як ядро для нейронної мережі та підбору "
            "гіперпараметрів. Особливістю моєї роботи є модель "
            "tuned_gam, яка буквально об'єднує обидва розділи в одну "
            "практичну реалізацію. Окрім того, при реалізації "
            "алгоритму з дисертації виявлено технічну проблему — "
            "колапс дисперсії при чистому EM-оновленні. Її розв'язано "
            "власною поправкою cov_lr, аналогічною rank-μ оновленню "
            "класичного CMA-ES. Весь код, документація та результати "
            "експериментів — у відкритому репозиторії GitHub. "
            "Дякую за увагу, готовий відповісти на запитання."
        ),
    },
]


# ============================================================================
# PPTX helpers

def _shape_rect(slide, x, y, w, h, fill, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def _text(slide, x, y, w, h, text, *, size=18, bold=False, color=None,
          align=None, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    # multi-line через перші параграфи
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    return tf


def _bullets(slide, x, y, w, h, bullets, *, size=16, color=None,
             marker_color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_top = Cm(0.05)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        # маркер
        m_run = p.add_run()
        m_run.text = "■  "
        m_run.font.name = "Calibri"
        m_run.font.size = Pt(size)
        m_run.font.color.rgb = marker_color or COLOR_ACCENT_1
        # сам текст
        t_run = p.add_run()
        t_run.text = line
        t_run.font.name = "Calibri"
        t_run.font.size = Pt(size)
        t_run.font.color.rgb = color or COLOR_TEXT


def _add_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    run = p.add_run()
    run.text = text


def _add_header(slide, prs, number, title):
    """Шапка слайду — кольорова смуга з номером і назвою."""
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, Cm(2.0),
                fill=COLOR_PRIMARY)
    _text(slide, Cm(0.7), Cm(0.2), Cm(2.5), Cm(1.6),
          number, size=30, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)
    _text(slide, Cm(3.0), Cm(0.45), Cm(20), Cm(1.4),
          title, size=24, bold=True,
          color=COLOR_WHITE, align=PP_ALIGN.LEFT)


def _add_footer(slide, prs):
    """Тонка акцентна лінія знизу слайду."""
    _shape_rect(slide, Cm(0), Cm(13.7), prs.slide_width, Cm(0.6),
                fill=COLOR_ACCENT_1)


# ============================================================================
# slide layouts

def _slide_title(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # лівий кольоровий бар
    _shape_rect(slide, Cm(0), Cm(0), Cm(1.5), prs.slide_height,
                fill=COLOR_PRIMARY)
    # заголовок
    _text(slide, Cm(2.5), Cm(4.5), Cm(22), Cm(3),
          data["title"], size=36, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    # лінія-розділювач
    _shape_rect(slide, Cm(10), Cm(8.5), Cm(5.4), Cm(0.1),
                fill=COLOR_ACCENT_2)
    # підзаголовок
    _text(slide, Cm(2.5), Cm(9.2), Cm(22), Cm(1.5),
          data["subtitle"], size=22,
          color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    # футер
    _text(slide, Cm(2.5), Cm(12.5), Cm(22), Cm(1.5),
          data["footer"], size=14,
          color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    _add_speaker_notes(slide, data["speech"])


def _slide_two_col(prs, data):
    """Шапка + ліворуч буллети, праворуч одна картинка."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, prs, data["number"], data["title"])

    # ліва колонка — буллети
    _bullets(slide, Cm(0.7), Cm(2.8), Cm(11.5), Cm(10.5),
             data["bullets"], size=16)

    # права колонка — рамка + картинка
    img_path = data.get("image")
    if img_path and Path(img_path).exists():
        _shape_rect(slide, Cm(12.7), Cm(2.6), Cm(12.0), Cm(10.7),
                    fill=COLOR_BG_LIGHT)
        slide.shapes.add_picture(str(img_path),
                                 Cm(13.0), Cm(2.9),
                                 width=Cm(11.4), height=Cm(10.1))

    _add_footer(slide, prs)
    _add_speaker_notes(slide, data["speech"])


def _slide_two_images(prs, data):
    """Шапка + буллети зверху, дві картинки знизу."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, prs, data["number"], data["title"])

    # буллети зверху
    _bullets(slide, Cm(0.7), Cm(2.5), Cm(24), Cm(3.5),
             data["bullets"], size=15)

    # дві картинки внизу
    for img, x in [(data.get("image_left"), Cm(0.7)),
                   (data.get("image_right"), Cm(12.9))]:
        if img and Path(img).exists():
            _shape_rect(slide, x, Cm(6.5), Cm(11.7), Cm(6.7),
                        fill=COLOR_BG_LIGHT)
            slide.shapes.add_picture(str(img), x + Cm(0.2), Cm(6.7),
                                     width=Cm(11.3), height=Cm(6.3))

    _add_footer(slide, prs)
    _add_speaker_notes(slide, data["speech"])


def _slide_big_image(prs, data):
    """Шапка + велика картинка зверху, буллети знизу."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, prs, data["number"], data["title"])

    # велика картинка
    img = data.get("image")
    if img and Path(img).exists():
        _shape_rect(slide, Cm(0.7), Cm(2.5), Cm(24), Cm(8),
                    fill=COLOR_BG_LIGHT)
        slide.shapes.add_picture(str(img), Cm(0.9), Cm(2.7),
                                 width=Cm(23.6), height=Cm(7.6))

    # підпис під картинкою
    if "caption" in data:
        _text(slide, Cm(0.7), Cm(10.6), Cm(24), Cm(0.6),
              data["caption"], size=11, color=COLOR_MUTED,
              align=PP_ALIGN.CENTER)

    # буллети знизу
    if "footer_bullets" in data:
        _bullets(slide, Cm(0.7), Cm(11.3), Cm(24), Cm(2.4),
                 data["footer_bullets"], size=13)

    _add_footer(slide, prs)
    _add_speaker_notes(slide, data["speech"])


# ============================================================================

def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.29)  # 16:9

    builders = {
        "title": _slide_title,
        "two_col": _slide_two_col,
        "two_images": _slide_two_images,
        "big_image": _slide_big_image,
    }
    for data in SLIDES:
        builders[data["kind"]](prs, data)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_PATH))
    return PPTX_PATH


# ============================================================================
# DOCX — текст для зачитування

def _docx_set_run(run, *, size=12, bold=False, italic=False,
                  font="Times New Roman"):
    run.font.name = font
    run.font.size = DPt(size)
    run.font.bold = bold
    run.font.italic = italic
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    rfonts.set(qn("w:ascii"), font)
    rfonts.set(qn("w:hAnsi"), font)
    rfonts.set(qn("w:cs"), font)


def _docx_para(doc, text, *, size=12, bold=False, italic=False,
               align=None, indent=DCm(1.25), space=DPt(6)):
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    pf = p.paragraph_format
    if indent is not None:
        pf.first_line_indent = indent
    pf.space_after = space
    pf.line_spacing = 1.5
    run = p.add_run(text)
    _docx_set_run(run, size=size, bold=bold, italic=italic)


def _docx_heading(doc, text, level=1):
    p = doc.add_paragraph()
    pf = p.paragraph_format
    pf.space_before = DPt(14)
    pf.space_after = DPt(8)
    pf.keep_with_next = True
    pf.line_spacing = 1.15
    p.alignment = DOCX_ALIGN.CENTER if level == 1 else DOCX_ALIGN.LEFT
    run = p.add_run(text)
    _docx_set_run(run, size=16 if level == 1 else 13, bold=True)


def build_docx() -> Path:
    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = DPt(12)
    for section in doc.sections:
        section.left_margin = DCm(2.5)
        section.right_margin = DCm(1.5)
        section.top_margin = DCm(2)
        section.bottom_margin = DCm(2)

    # титулка
    for _ in range(4):
        doc.add_paragraph()
    _docx_para(doc, "ТЕКСТ УСНОГО ВИСТУПУ",
               size=20, bold=True, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(8))
    _docx_para(doc, "Презентація на 6 слайдів (3-5 хвилин)",
               size=14, italic=True, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(20))
    _docx_para(doc, "Порівняння методів класифікації\n"
                    "з розширеним алгоритмом CMA-ES",
               size=14, bold=True, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(60))
    _docx_para(doc, "Богдан Волошенюк",
               size=12, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(2))
    _docx_para(doc, "ЧНУ ім. Юрія Федьковича, 2026",
               size=12, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(20))
    doc.add_page_break()

    # підказка
    _docx_heading(doc, "Як користуватися цим текстом", level=1)
    _docx_para(doc,
        "Цей файл містить повний текст для усного виступу. Тривалість "
        "приблизно 3-5 хвилин при спокійному темпі мовлення (близько "
        "130 слів на хвилину). На кожний слайд відведено 30-50 секунд. "
        "Перед виступом рекомендується одноразово прочитати вголос, "
        "щоб відчути ритм і виправити складні слова.")
    _docx_para(doc,
        "Текст можна тримати перед собою на захисті як підказку. "
        "Сама ж презентація показується великому екрану. Якщо "
        "користуватися Режимом доповідача в PowerPoint (Presenter "
        "View), той самий текст видно у нотатках до кожного слайда.")
    doc.add_page_break()

    # суцільний текст всієї доповіді
    _docx_heading(doc, "Суцільний текст доповіді", level=1)
    _docx_para(doc,
        "Нижче — повний текст для зачитування поспіль. Це версія для "
        "того, хто хоче читати без перемикання між слайдами в "
        "комп'ютерному поданні. На захисті між абзацами доцільно "
        "робити паузи 2-3 секунди для переходу на наступний слайд.",
        italic=True)
    for i, data in enumerate(SLIDES, start=1):
        _docx_para(doc, data["speech"], size=12)
    doc.add_page_break()

    # текст по слайдах
    _docx_heading(doc, "Текст по слайдах", level=1)
    for i, data in enumerate(SLIDES, start=1):
        _docx_heading(doc, f"Слайд {i}. {data['title'].replace(chr(10), ' ')}",
                      level=2)
        if "bullets" in data:
            _docx_para(doc, "Що показано на слайді:",
                       bold=True, indent=None, space=DPt(2))
            for b in data["bullets"]:
                p = doc.add_paragraph()
                pf = p.paragraph_format
                pf.left_indent = DCm(0.75)
                pf.first_line_indent = DCm(-0.5)
                pf.space_after = DPt(3)
                pf.line_spacing = 1.4
                run = p.add_run("•  " + b)
                _docx_set_run(run, size=11)
        elif "footer_bullets" in data:
            _docx_para(doc,
                       f"Зображення: {Path(data['image']).name}",
                       italic=True, indent=None, space=DPt(2))
            _docx_para(doc, "Основні висновки:",
                       bold=True, indent=None, space=DPt(2))
            for b in data["footer_bullets"]:
                p = doc.add_paragraph()
                pf = p.paragraph_format
                pf.left_indent = DCm(0.75)
                pf.first_line_indent = DCm(-0.5)
                pf.space_after = DPt(3)
                pf.line_spacing = 1.4
                run = p.add_run("•  " + b)
                _docx_set_run(run, size=11)

        _docx_para(doc, "Що читати:",
                   bold=True, indent=None, space=DPt(2))
        _docx_para(doc, data["speech"], size=12)
        doc.add_paragraph()

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(DOCX_PATH))
    return DOCX_PATH


def main() -> int:
    print("Генерую PowerPoint...")
    pptx_out = build_pptx()
    print(f"  OK: {pptx_out} ({pptx_out.stat().st_size // 1024} KB)")

    print("Генерую текст доповіді (.docx)...")
    docx_out = build_docx()
    print(f"  OK: {docx_out} ({docx_out.stat().st_size // 1024} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())

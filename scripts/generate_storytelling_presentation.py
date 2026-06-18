"""Сторітелінгова презентація – 12 слайдів за 5-актною структурою.

HOOK (1-2)  → PROBLEM (3-4) → INSIGHT (5-7) → SOLUTION (8-10) → CONCLUSION (11-12)

Цільова аудиторія: екзаменаційна комісія кафедри прикладної математики
ЧНУ. Тривалість виступу – 8-10 хвилин. Прізвище Юлії Анатоліївни не
згадується (за вимогою). У кінці – QR-код на GitHub.

Файли:
* ``../../presentation/storytelling_presentation.pptx``
* ``../../presentation/storytelling_presentation_speech.docx``
"""

from __future__ import annotations

import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH as DOCX_ALIGN
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Cm as DCm, Pt as DPt

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt, Emu


REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = REPO_ROOT.parent / "presentation"
PPTX_PATH = OUT_DIR / "storytelling_presentation.pptx"
DOCX_PATH = OUT_DIR / "storytelling_presentation_speech.docx"

DOCS_FIG = REPO_ROOT.parent / "docs" / "figures"
DIPLOMA_FIG = REPO_ROOT.parent / "diploma_figures"


# Палітра (тепла, академічна)
COLOR_PRIMARY = RGBColor(0x26, 0x46, 0x53)   # глибокий бірюзовий
COLOR_ACCENT_1 = RGBColor(0x2A, 0x9D, 0x8F)  # м'ятний
COLOR_ACCENT_2 = RGBColor(0xE7, 0x6F, 0x51)  # кораловий
COLOR_ACCENT_3 = RGBColor(0xE9, 0xC4, 0x6A)  # золотий
COLOR_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_MUTED = RGBColor(0x6C, 0x6C, 0x6C)
COLOR_BG_LIGHT = RGBColor(0xFD, 0xFC, 0xFA)
COLOR_BG_DARK = RGBColor(0x1F, 0x33, 0x3D)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SHADOW = RGBColor(0xD0, 0xD0, 0xD0)


# =============================================================================
# контент 12 слайдів за актами

SLIDES = [
    # ============================== ACT 1: HOOK ==============================
    # ---- Слайд 1 — Титульний з керівником ----
    {
        "kind": "title_hero",
        "act": "HOOK",
        # U+2011 NON-BREAKING HYPHEN, щоб PowerPoint не розривав "CMA-ES"
        "title": "CMA‑ES",
        "subtitle": ("як інструмент налаштування параметрів\n"
                     "в алгоритмах класифікації\n\n"
                     "Випускна кваліфікаційна робота бакалавра\n"
                     "Спеціальність 122 «Комп'ютерні науки»"),
        "author_block": [
            "Виконав: Волошенюк Богдан Анатолійович",
            "група 441-А",
            "Керівник: д. фіз.-мат. наук, проф. Малик І. В.",
            "ЧНУ ім. Юрія Федьковича · 2026",
        ],
        "speech": (
            "Шановний голово, шановні члени комісії. Представляю "
            "кваліфікаційну роботу за темою «CMA-ES як інструмент "
            "налаштування параметрів в алгоритмах класифікації». "
            "Виконав Волошенюк Богдан Анатолійович, група 441-А, "
            "спеціальність 122 – Комп'ютерні науки. Науковий керівник – "
            "кандидат фізико-математичних наук, доцент Малик Ігор "
            "Володимирович."
        ),
    },

    # ---- Слайд 2 – числа, що чіпляють ----
    {
        "kind": "big_numbers",
        "act": "HOOK",
        "label": "Масштаб дослідження",
        "numbers": [
            ("3", "датасети", "UCI ML Repository", COLOR_ACCENT_2),
            ("12", "моделей", "порівняно", COLOR_ACCENT_1),
            ("34", "тести", "усі проходять", COLOR_ACCENT_3),
        ],
        "tagline": "Усі цифри – з реальних прогонів програми",
        "speech": (
            "Спочатку про масштаб. Дослідження проведено на трьох сучасних "
            "відкритих датасетах з UCI Machine Learning Repository. "
            "Реалізовано 12 моделей класифікації – від класичної "
            "логістичної регресії до нейронних мереж, навчених без "
            "використання градієнтів. Покрито 34 автоматизованими тестами "
            "на pytest, усі проходять. Усі цифри – з реальних прогонів "
            "програми, не симуляція."
        ),
    },

    # ============================== ACT 2: PROBLEM ===========================
    # ---- Слайд 3 – головна проблема ML ----
    {
        "kind": "problem_question",
        "act": "PROBLEM",
        "label": "Проблема",
        "question": "Який метод класифікації\nобрати для нової задачі?",
        "options": [
            "Логістична регресія?",
            "SVM з RBF-ядром?",
            "Нейронна мережа?",
            "Або щось зовсім інше?",
        ],
        "punchline": "Відповідь залежить від конкретної задачі.\nУніверсального найкращого методу – НЕ існує.",
        "speech": (
            "Уявіть собі ситуацію. Перед вами нова задача класифікації. "
            "Яку модель обрати? Логістичну регресію – швидко, але "
            "лінійно. SVM – потужно, але повільно. Нейронну мережу – "
            "гнучко, але незрозуміло чому. Сучасне машинне навчання "
            "виробило важливий принцип: універсального найкращого методу "
            "не існує. Вибір завжди залежить від характеру конкретної "
            "задачі. Це і є перша велика проблема, з якою стикаються "
            "інженери даних."
        ),
    },

    # ---- Слайд 4 – друга проблема: гіперпараметри ----
    {
        "kind": "problem_split",
        "act": "PROBLEM",
        "label": "Проблема",
        "title": "Друга проблема – гіперпараметри",
        "left_title": "Що це?",
        "left_text": (
            "Параметри САМОЇ МОДЕЛІ, не вагів.\n\n"
            "Задаються ДО навчання.\n\n"
            "Контролюють поведінку алгоритму."
        ),
        "right_title": "Чому це складно?",
        "right_text": (
            "Простір параметрів – безперервний.\n\n"
            "Параметри взаємопов'язані.\n\n"
            "Перебір вручну – не варіант."
        ),
        "footer": "Потрібен інтелектуальний автоматичний підбір.",
        "speech": (
            "Друга велика проблема – підбір гіперпараметрів. Це не вагові "
            "коефіцієнти моделі, а її параметри. Скільки сусідів брати у "
            "kNN? Яка ширина ядра у SVM? Скільки нейронів у прихованому "
            "шарі? Усі ці значення треба задати ДО початку навчання. "
            "І якщо їх вибрати неправильно – модель буде працювати "
            "погано, незалежно від того, наскільки хороший базовий "
            "алгоритм. Простір параметрів зазвичай безперервний і "
            "багатовимірний, параметри взаємопов'язані. Перебір вручну "
            "займає тижні. Потрібен інтелектуальний автоматичний підбір."
        ),
    },

    # ============================== ACT 3: INSIGHT ===========================
    # ---- Слайд 5 – GAM ----
    {
        "kind": "two_col_card",
        "act": "INSIGHT",
        "label": "Розділ 1 теорії",
        "title": "GAM – гнучкість без жертв",
        "left_bullets": [
            "Узагальнена адитивна модель",
            "Розширення лінійної регресії",
            "Гладкі функції для кожної ознаки",
            "Сума функцій впливу",
            "B-сплайнові базиси",
        ],
        "right_quote": (
            "Гнучкість нейронної мережі\n"
            "+\n"
            "інтерпретованість регресії"
        ),
        "speech": (
            "Перший теоретичний розділ присвячено узагальненим адитивним "
            "моделям. Це красиве рішення на стику двох світів. З одного "
            "боку – гнучкість нейронних мереж, що вловлюють нелінійні "
            "залежності. З іншого – інтерпретованість лінійних моделей: "
            "можна окремо побачити вплив кожної ознаки. Технічно GAM "
            "розкладає кожну ознаку через B-сплайнові базисні функції, "
            "потім лінійно комбінує їх через логістичну регресію. У "
            "результаті – модель, яка одночасно гнучка і прозора."
        ),
    },

    # ---- Слайд 6 – CMA-ES ----
    {
        "kind": "two_col_card",
        "act": "INSIGHT",
        "label": "Розділ 2 теорії",
        "title": "CMA-ES – оптимізатор без похідних",
        "left_bullets": [
            "Еволюційна стратегія",
            "Стохастичний пошук",
            "Адаптує область пошуку",
            "Не потребує градієнтів",
            "Один з найкращих у класі",
        ],
        "right_quote": (
            "Працює там,\n"
            "де класичні методи\n"
            "не справляються"
        ),
        "speech": (
            "Другий теоретичний розділ – алгоритм адаптації коваріаційної "
            "матриці. Це еволюційна стратегія, яка не потребує "
            "обчислення похідних цільової функції. Алгоритм самостійно "
            "адаптує область пошуку, враховує кореляції між параметрами, "
            "автоматично балансує між дослідженням нових регіонів і "
            "уточненням знайдених рішень. У задачах оптимізації чорної "
            "скриньки CMA-ES вважається одним з найкращих методів."
        ),
    },

    # ---- Слайд 7 – Розширений CMA-ES (Літвінчук, я розширив) ----
    {
        "kind": "two_col_card",
        "act": "INSIGHT",
        "label": "Розвиток теорії",
        "title": "Розширений CMA-ES зі сумішами розподілів",
        "left_bullets": [
            "Унімодальний розподіл → суміш k компонент",
            "Параметри суміші оцінюються EM-алгоритмом",
            "Самоадаптивний підбір кількості піків",
            "Теоретичні переваги – на мультимодальних задачах",
            "Дослідження проводила Літвінчук Ю. А. (2024)",
        ],
        "right_quote": (
            "Класичний підхід – Літвінчук.\n"
            "Програмна реалізація\n"
            "і розширення – мої."
        ),
        "speech": (
            "Дослідження розширеного варіанту CMA-ES зі сумішами "
            "нормальних розподілів проводила Літвінчук Юлія Анатоліївна "
            "у дисертаційній роботі 2024 року. У моїй роботі я взяв цей "
            "теоретичний підхід і розширив його у програмному напрямку: "
            "самостійно реалізував алгоритм з нуля, додав механізм "
            "самоадаптивного підбору кількості піків та технічне "
            "доповнення для стабільної збіжності на унімодальних задачах. "
            "Параметри суміші – ваги, центри і коваріації – оновлюються "
            "EM-алгоритмом за найкращою половиною хромосом кожної "
            "ітерації."
        ),
    },

    # ---- Слайд 7 – АГА-МОМЕНТ ----
    {
        "kind": "insight_dark",
        "act": "INSIGHT",
        "label": "Ключова ідея",
        "big_quote": (
            "А що, якщо взяти алгоритм з другого розділу\n"
            "і застосувати його до моделі з першого?"
        ),
        "subtitle_quote": (
            "CMA-ES шукає оптимальні параметри GAM –\n"
            "теорія об'єднується у працюючу модель"
        ),
        "model_name": "tuned_gam",
        "speech": (
            "Тепер головна ідея всієї роботи. Два теоретичні розділи "
            "диплома – GAM і CMA-ES – на перший погляд не мають "
            "перетину. Один про моделі класифікації, інший про "
            "оптимізатори. Але... що, якщо взяти алгоритм з другого "
            "розділу і застосувати його до моделі з першого? Що, якщо "
            "CMA-ES шукатиме оптимальні параметри для GAM? У результаті "
            "теорія перестає бути просто описом – вона стає працюючою "
            "моделлю tuned_gam, яка демонструє єдність обох "
            "теоретичних концепцій."
        ),
    },

    # ============================== ACT 4: SOLUTION ==========================
    # ---- Слайд 9 – Порівняльний аналіз: задача ----
    {
        "kind": "two_col_card",
        "act": "SOLUTION",
        "label": "Порівняльний аналіз",
        "title": "Що саме порівнюємо",
        "left_bullets": [
            "12 моделей у трьох групах: 5 базових,",
            "2 CMA-NN, 5 із tuned-підбором гіперпараметрів",
            "3 датасети UCI ML Repository різного типу",
            "Єдині умови оцінки для всіх моделей",
            "3 класичні метрики: Accuracy, F1, ROC-AUC",
        ],
        "right_quote": (
            "Чесне порівняння –\n"
            "однакові умови\n"
            "для всіх моделей"
        ),
        "speech": (
            "Перш ніж переходити до результатів, кілька слів про "
            "методологію порівняльного аналізу. Я зіставив дванадцять "
            "моделей у трьох групах: п'ять базових алгоритмів машинного "
            "навчання, два CMA-NN з безградієнтним навчанням ваг та "
            "п'ять tuned-моделей з автоматичним підбором гіперпараметрів. "
            "Усі моделі тестувалися на трьох сучасних відкритих "
            "датасетах в однакових умовах – та сама перехресна "
            "валідація, ті самі метрики, той самий random seed. Це "
            "забезпечує чесність порівняння."
        ),
    },

    # ---- Слайд 10 – Порівняльний аналіз: підхід ----
    {
        "kind": "method_pipeline",
        "act": "SOLUTION",
        "label": "Порівняльний аналіз",
        "title": "Конвеєр порівняння",
        "pipeline_steps": [
            ("01", "Препроцесинг", "Один на всі моделі"),
            ("02", "Навчання", "12 моделей паралельно"),
            ("03", "Оцінка", "5-fold CV, 3 метрики"),
            ("04", "Ранжування", "Лідери на кожному датасеті"),
        ],
        "tagline": "Жоден метод не перемагає на всіх задачах – принцип NFL",
        "speech": (
            "Конвеєр порівняння складається з чотирьох кроків. Спочатку "
            "однаковий препроцесинг для всіх моделей. Далі – паралельне "
            "навчання дванадцяти моделей на одних і тих самих даних. "
            "Потім – оцінка через п'ятикратну перехресну валідацію за "
            "трьома метриками. На фінальному кроці – ранжування моделей "
            "і визначення лідерів окремо для кожного датасета. Ключовий "
            "висновок дослідження – жоден метод не перемагає одночасно "
            "на всіх задачах, що повністю узгоджується з принципом "
            "«немає безплатних обідів»."
        ),
    },

    # ---- Слайд 11 – Архітектура ----
    {
        "kind": "architecture_slide",
        "act": "SOLUTION",
        "label": "Реалізація",
        "title": "Як це працює зсередини",
        "image": DOCS_FIG / "01_architecture.png",
        "right_stats": [
            ("Python", "мова"),
            ("13", "модулів"),
            ("≈2500", "рядків"),
            ("34", "тести"),
        ],
        "speech": (
            "Програмна реалізація побудована як модульний Python-пакет. "
            "Кожен модуль відповідає за одну конкретну задачу. Це "
            "свідоме архітектурне рішення: такий код легше тестувати, "
            "легше додавати нові моделі або датасети, а новий розробник "
            "за п'ять хвилин розуміє, де що лежить. У цифрах: "
            "тринадцять модулів, близько двох з половиною тисяч рядків "
            "коду, тридцять чотири тести на pytest, усі проходять "
            "успішно."
        ),
    },

    # ---- Слайд 9 – Методика ----
    {
        "kind": "method_pipeline",
        "act": "SOLUTION",
        "label": "Методика",
        "title": "Як ми міряємо якість",
        "pipeline_steps": [
            ("01", "Train/Test 80:20", "Стратифікований поділ"),
            ("02", "5-fold CV", "Перевірка стабільності"),
            ("03", "3 метрики", "Accuracy · F1 · AUC"),
            ("04", "MLflow", "Відтворюваність"),
        ],
        "tagline": "Стандартний sklearn API + журналювання експериментів",
        "speech": (
            "Методика оцінки – стандартна для машинного навчання, без "
            "сюрпризів. Спочатку дані діляться на тренувальну і тестову "
            "вибірки у пропорції вісімдесят на двадцять, зі стратифікацією. "
            "На тренувальній частині додатково проводиться п'ятикратна "
            "перехресна валідація. Обчислюються три метрики, затверджені "
            "куратором: accuracy, F1-score, ROC-AUC. Усі експерименти "
            "журналюються через MLflow для повної відтворюваності."
        ),
    },

    # ---- Слайд: Таблиця метрик (PhiUSIIL) ----
    {
        "kind": "image_with_caption",
        "act": "SOLUTION",
        "label": "Порівняння метрик",
        "title": "Результати на датасеті PhiUSIIL",
        "image": DIPLOMA_FIG / "metrics_table_phiusiil.png",
        "tagline": ("Усі 12 моделей у єдиній таблиці: Accuracy, F1, AUC "
                    "та час навчання. GAM і tuned_gam – ідеальні 1.0000."),
        "speech": (
            "Перш ніж перейти до підсумків, покажу повну таблицю метрик "
            "на найпростішому датасеті – PhiUSIIL. У ній зведено "
            "дванадцять моделей з трьома стандартними метриками: "
            "Accuracy, F1-score та ROC-AUC, плюс час навчання. Відразу "
            "видно безумовних лідерів – GAM і tuned_gam досягли "
            "ідеальних одиниць по всіх метриках. Майже всі моделі "
            "показали якість вище 0,99 – задача в значній мірі лінійно "
            "роздільна. Цікаво порівняти час: класична логістична "
            "регресія – одна сота секунди, тоді як cma_mixture – понад "
            "двадцять секунд. Це наочно показує плату за безградієнтну "
            "оптимізацію."
        ),
    },

    # ---- Слайд: PhiUSIIL детально ----
    {
        "kind": "two_col_card",
        "act": "SOLUTION",
        "label": "Результати – датасет №1",
        "title": "PhiUSIIL Phishing URL",
        "left_bullets": [
            "Бінарна, ≈3 000 рядків, 54 ознаки URL",
            "Сучасні дані – публікація 2024 року",
            "Майже лінійно роздільна задача",
            "Безумовний лідер – GAM та tuned_gam",
            "F1 = 1.0000, AUC = 1.0000",
        ],
        "right_quote": (
            "Ідеальні метрики –\n"
            "0 помилок\n"
            "на 600 прикладах"
        ),
        "speech": (
            "Перший датасет – PhiUSIIL: класифікація phishing-сайтів за "
            "ознаками URL-адреси. Дані за 2024 рік, актуальні. Задача "
            "виявилася майже лінійно роздільною. Безумовний лідер – GAM "
            "та tuned_gam з ідеальними значеннями всіх метрик: на шести "
            "сотнях тестових прикладів жодної помилки. Це підтверджує "
            "силу B-сплайнового базисного перетворення для нелінійно "
            "роздільних задач."
        ),
    },

    # ---- Слайд: Таблиця метрик (Steel Plate Defects) ----
    {
        "kind": "image_with_caption",
        "act": "SOLUTION",
        "label": "Порівняння метрик",
        "title": "Результати на датасеті Steel Plate Defects",
        "image": DIPLOMA_FIG / "metrics_table_steel_plate.png",
        "tagline": ("Мультикласова задача – 7 типів дефектів. "
                    "Лідер: tuned_svm з F1=0.7597, AUC=0.9281."),
        "speech": (
            "Друга таблиця – Steel Plate Defects, мультикласова задача "
            "з семи типів дефектів металевих пластин. Метрики тут "
            "очікувано нижчі за бінарну задачу. Безумовний лідер – "
            "tuned_svm з F1 0.7597 та AUC 0.9281; покращення відносно "
            "базового SVM становить майже сотку. Слідом – tuned_gam з "
            "F1 0.7576. Безградієнтні моделі CMA-NN показали слабше: "
            "для семи класів і двохсот ваг тридцяти ітерацій "
            "недостатньо. Час навчання cma_mixture – близько 12 секунд."
        ),
    },

    # ---- Слайд 15 – Steel Plate детально ----
    {
        "kind": "two_col_card",
        "act": "SOLUTION",
        "label": "Результати – датасет №2",
        "title": "Steel Plate Defects",
        "left_bullets": [
            "Мультиклас, 7 типів дефектів",
            "27 ознак, 1 941 рядок",
            "Лідер – tuned_svm: F1 = 0.7597, AUC = 0.9281",
            "Покращення +0.0083 над базовим SVM",
            "CMA-NN: 30 ітерацій недостатньо для 200 ваг",
        ],
        "right_quote": (
            "Найкраща модель –\n"
            "автоматичний підбір C і γ\n"
            "через CMA-ES"
        ),
        "speech": (
            "Другий датасет – промислова класифікація дефектів металу. "
            "Сім класів, двадцять сім фізичних ознак пластини. Лідером "
            "став tuned_svm з F1 близько 0.76. Покращення відносно "
            "базового SVM – вісім тисячних, що для мультикласу істотно. "
            "Безградієнтні моделі CMA-NN показали слабше: для семи "
            "класів і двохсот ваг тридцяти ітерацій недостатньо для "
            "збіжності. Це обмеження безградієнтного підходу на "
            "великорозмірних задачах."
        ),
    },

    # ---- Слайд: Таблиця метрик (Credit Default) ----
    {
        "kind": "image_with_caption",
        "act": "SOLUTION",
        "label": "Порівняння метрик",
        "title": "Результати на датасеті Default of Credit Card Clients",
        "image": DIPLOMA_FIG / "metrics_table_credit_default.png",
        "tagline": ("Незбалансовані 78/22. Лідер: tuned_mlp (F1=0.4360). "
                    "Розширений CMA-ES вперше випереджає класичний."),
        "speech": (
            "Третя таблиця – Default of Credit Card Clients. Найскладніша "
            "задача у дослідженні через сильну незбалансованість класів "
            "сімдесят вісім на двадцять два. Лідер за F1 – tuned_mlp з "
            "0.4360, поряд tuned_svm (0.4356) і tuned_gam (0.4322). "
            "Найвищий AUC – у базового GAM, 0.7311. Ключовий і "
            "найнесподіваніший результат: cma_mixture істотно "
            "випередив cma_classic за F1: 0.3333 проти 0.2000. Це "
            "перше у дослідженні емпіричне підтвердження теоретичної "
            "переваги розширеного варіанту CMA-ES."
        ),
    },

    # ---- Слайд 16 – Credit Default + ключовий результат ----
    {
        "kind": "two_col_card",
        "act": "SOLUTION",
        "label": "Результати – датасет №3",
        "title": "Default of Credit Card Clients",
        "left_bullets": [
            "Бінарна, 30 000 клієнтів, 23 ознаки",
            "Сильна незбалансованість 78/22",
            "Лідер – tuned_mlp: F1 = 0.4360",
            "GAM має найвищий AUC = 0.7311",
            "cma_mixture > cma_classic: F1 0.33 vs 0.20",
        ],
        "right_quote": (
            "Тут розширений CMA-ES\n"
            "вперше випередив\n"
            "класичний варіант"
        ),
        "speech": (
            "Третій датасет – прогнозування дефолту клієнтів кредитних "
            "карток. Сильна незбалансованість сімдесят вісім на двадцять "
            "два робить задачу складною. Лідер серед усіх моделей – "
            "tuned_mlp. Однак найцікавіший результат – на цій задачі "
            "розширений варіант CMA-ES зі сумішами вперше у дослідженні "
            "випередив класичний: F1 нуль ціла тридцять три проти двох "
            "десятих. Це емпіричне підтвердження теоретичної гіпотези "
            "про перевагу сумішевого варіанту на мультимодальних "
            "ландшафтах цільової функції."
        ),
    },

    # ---- Слайд 17 – GUI Streamlit ----
    {
        "kind": "two_col_card",
        "act": "SOLUTION",
        "label": "Інтерфейс",
        "title": "Дашборд на Streamlit",
        "left_bullets": [
            "Інтерактивний веб-інтерфейс одним кліком",
            "Вибір датасета, моделей, параметрів",
            "Прогрес-бар і лог у реальному часі",
            "Таблиця метрик з адаптивними progress-барами",
            "Експорт у CSV, інтеграція з MLflow",
        ],
        "right_quote": (
            "Жива демонстрація\n"
            "доступна для перегляду\n"
            "у будь-який момент"
        ),
        "speech": (
            "Окремо реалізовано графічний веб-дашборд на Streamlit. У "
            "боковій панелі користувач послідовно обирає всі параметри "
            "запуску. Після старту з'являється прогрес-бар і лог у "
            "реальному часі. Результати виводяться у вигляді "
            "інтерактивної таблиці з кольоровими progress-барами на "
            "адаптивній шкалі, що дозволяє візуально бачити різницю "
            "навіть між дуже близькими значеннями. Усі експерименти "
            "автоматично логуються через MLflow."
        ),
    },

    # ============================== ACT 5: CONCLUSION ========================
    # ---- Слайд 11 – Технічна деталь реалізації ----
    {
        "kind": "contribution",
        "act": "CONCLUSION",
        "label": "Деталь реалізації",
        "title": "Параметр cov_lr",
        "story_block": (
            "У тестах на унімодальних функціях\n"
            "чисте EM-оновлення коваріаційних матриць\n"
            "виходило нестабільним за 5–7 ітерацій."
        ),
        "fix_label": "Параметр cov_lr",
        "fix_text": (
            "Момент-усереднення коваріацій між ітераціями.\n"
            "За аналогією з rank-µ оновленням класичного CMA-ES.\n"
            "За замовчуванням 0,2."
        ),
        "result_text": "Алгоритм сходиться на унімодальних і мультимодальних задачах.",
        "speech": (
            "Невелика технічна деталь з практики. У тестах на простих "
            "унімодальних функціях помітив, що чисте EM-оновлення "
            "коваріаційних матриць давало нестабільну поведінку: "
            "дисперсія швидко зменшувалася і пошук зупинявся вже після "
            "п'яти–семи ітерацій. Щоб обійти цю особливість, додав у "
            "реалізацію параметр cov_lr – коефіцієнт момент-усереднення "
            "коваріацій між ітераціями, за аналогією з rank-µ оновленням "
            "у класичному CMA-ES. За замовчуванням використовується "
            "значення 0,2. З цією поправкою алгоритм стабільно сходиться "
            "як на унімодальних, так і на мультимодальних задачах."
        ),
    },

    # ---- Слайд 19 – Висновки ----
    {
        "kind": "two_col_card",
        "act": "CONCLUSION",
        "label": "Висновки",
        "title": "Основні результати роботи",
        "left_bullets": [
            "Реалізовано та порівняно 12 моделей класифікації",
            "GAM-класифікатор – у лідерах на 2 з 3 датасетах",
            "tuned_* стабільно покращують базові моделі",
            "Розширений CMA-ES емпірично перевершує класичний",
            "Підтверджено принцип «немає безплатних обідів»",
        ],
        "right_quote": (
            "Мету досягнуто –\n"
            "усі завдання\n"
            "виконано"
        ),
        "speech": (
            "Підсумовуючи. Поставлену мету досягнуто, всі завдання "
            "виконано. Реалізовано і порівняно дванадцять моделей "
            "класифікації. Власна реалізація GAM-класифікатора увійшла "
            "до групи лідерів на двох з трьох датасетах. Підбір "
            "гіперпараметрів через CMA-ES стабільно покращує базові "
            "моделі. Розширений варіант CMA-ES зі сумішами емпірично "
            "перевершив класичний на незбалансованій задачі – перше "
            "практичне підтвердження теоретичної гіпотези. Жоден метод "
            "не переважає на всіх задачах – це повністю узгоджується з "
            "принципом «немає безплатних обідів»."
        ),
    },

    # ---- Слайд 20 – Дякую за увагу (велике посередині) ----
    {
        "kind": "thanks_big",
        "act": "CONCLUSION",
        "title": "Дякую за увагу!",
        "subtitle": "Готовий відповісти на ваші запитання",
        "github_label": "Повний код проєкту:",
        "github_url": "github.com/bohdanvolosheniuk1/classification_cma_es",
        "qr_image": DIPLOMA_FIG / "qr_github.png",
        "speech": (
            "Дякую за увагу. Готовий відповісти на ваші запитання."
        ),
    },
]


# =============================================================================
# PPTX helpers

def _shape_rect(slide, x, y, w, h, fill, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def _shape_rounded(slide, x, y, w, h, fill, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def _text(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
          color=None, align=None, font="Calibri", anchor=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
    return tf


def _bullet_list(slide, x, y, w, h, items, *, size=14, color=None,
                 marker_color=None, marker="■"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        m = p.add_run()
        m.text = f"{marker}  "
        m.font.name = "Calibri"
        m.font.size = Pt(size)
        m.font.color.rgb = marker_color or COLOR_ACCENT_1
        t = p.add_run()
        t.text = line
        t.font.name = "Calibri"
        t.font.size = Pt(size)
        t.font.color.rgb = color or COLOR_TEXT


def _act_badge(slide, prs, label):
    """Маленька плашка зверху правого кута з номером акту/міткою."""
    w = Cm(5)
    x = prs.slide_width - w - Cm(0.8)
    _shape_rounded(slide, x, Cm(0.5), w, Cm(0.9),
                   fill=COLOR_ACCENT_2)
    _text(slide, x, Cm(0.5), w, Cm(0.9),
          label, size=13, bold=True, color=COLOR_WHITE,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _footer_brand(slide, prs):
    """Мала акцентна смуга знизу."""
    _shape_rect(slide, Cm(0), prs.slide_height - Cm(0.4),
                prs.slide_width, Cm(0.4),
                fill=COLOR_ACCENT_1)


def _add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    p.add_run().text = text


def _shadow_image(slide, img_path, x, y, w, h, shadow_offset=Cm(0.2)):
    """Картинка з псевдо-тінню (сірий зсунутий прямокутник позаду)."""
    _shape_rect(slide, x + shadow_offset, y + shadow_offset, w, h,
                fill=COLOR_SHADOW)
    slide.shapes.add_picture(str(img_path), x, y, width=w, height=h)


# =============================================================================
# SLIDE 1 – title hero

def _slide_title_hero(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # темна ліва панель
    _shape_rect(slide, Cm(0), Cm(0), Cm(9.5), prs.slide_height,
                fill=COLOR_PRIMARY)
    # 3 акцентні смуги у нижньому лівому куті
    _shape_rect(slide, Cm(0), prs.slide_height - Cm(1.0),
                Cm(3.2), Cm(0.25), fill=COLOR_ACCENT_2)
    _shape_rect(slide, Cm(0), prs.slide_height - Cm(0.7),
                Cm(2.0), Cm(0.25), fill=COLOR_ACCENT_1)
    _shape_rect(slide, Cm(0), prs.slide_height - Cm(0.4),
                Cm(1.0), Cm(0.25), fill=COLOR_ACCENT_3)

    # величезний заголовок на темному фоні зліва (одне слово CMA‑ES)
    _text(slide, Cm(0.2), Cm(4.5), Cm(9.1), Cm(5),
          data["title"], size=66, bold=True,
          color=COLOR_WHITE, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)

    # підзаголовок справа (5 рядків, потрібен більший блок)
    _text(slide, Cm(10), Cm(2.0), Cm(14.5), Cm(5.5),
          data["subtitle"], size=18, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

    # автор – компактний блок з кольоровою смугою зліва
    _shape_rect(slide, Cm(10), Cm(8.3), Cm(0.2), Cm(4),
                fill=COLOR_ACCENT_2)
    for i, line in enumerate(data["author_block"]):
        sz = 14 if i == 0 else 12
        bold = (i == 0)
        col = COLOR_TEXT if i == 0 else COLOR_MUTED
        _text(slide, Cm(10.5), Cm(8.3 + i * 0.85), Cm(14.5), Cm(0.9),
              line, size=sz, bold=bold, color=col, align=PP_ALIGN.LEFT)

    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 2 – big numbers

def _slide_big_numbers(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # фон світлий
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    # маленький лейбл вгорі
    _text(slide, Cm(1), Cm(1.0), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)

    # три великі цифри в карточках по центру
    n_cards = len(data["numbers"])
    card_w = Cm(7.0)
    gap = Cm(0.7)
    total_w = card_w * n_cards + gap * (n_cards - 1)
    start_x = (prs.slide_width - total_w) // 2
    for i, (num, label_top, label_bot, color) in enumerate(data["numbers"]):
        x = start_x + (card_w + gap) * i
        # тінь
        _shape_rounded(slide, x + Cm(0.2), Cm(3.4), card_w, Cm(7.5),
                       fill=COLOR_SHADOW)
        # картка
        _shape_rounded(slide, x, Cm(3.2), card_w, Cm(7.5),
                       fill=COLOR_WHITE)
        # кольорова верхня смуга
        _shape_rect(slide, x, Cm(3.2), card_w, Cm(0.4),
                    fill=color)
        # величезне число
        _text(slide, x, Cm(4.0), card_w, Cm(3.5),
              num, size=72, bold=True, color=color,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # підпис під числом
        _text(slide, x, Cm(7.7), card_w, Cm(1),
              label_top, size=16, bold=True, color=COLOR_PRIMARY,
              align=PP_ALIGN.CENTER)
        _text(slide, x, Cm(8.7), card_w, Cm(1),
              label_bot, size=12, color=COLOR_MUTED,
              align=PP_ALIGN.CENTER)

    # tagline знизу
    _text(slide, Cm(1), Cm(12.2), prs.slide_width - Cm(2), Cm(0.8),
          data["tagline"], size=13, italic=True,
          color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 3 – provocative question

def _slide_problem_question(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    _text(slide, Cm(1), Cm(1.0), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)

    # велике питання вгорі по центру
    _text(slide, Cm(2), Cm(2.5), Cm(21.4), Cm(3.5),
          data["question"], size=38, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)

    # ряд варіантів – як кнопки-картки
    n = len(data["options"])
    btn_w = Cm(5.3)
    btn_h = Cm(1.8)
    gap = Cm(0.4)
    total_w = btn_w * n + gap * (n - 1)
    start_x = (prs.slide_width - total_w) // 2
    for i, opt in enumerate(data["options"]):
        x = start_x + (btn_w + gap) * i
        _shape_rounded(slide, x, Cm(6.8), btn_w, btn_h,
                       fill=COLOR_WHITE,
                       line_color=COLOR_PRIMARY)
        _text(slide, x, Cm(6.8), btn_w, btn_h,
              opt, size=14, color=COLOR_PRIMARY,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # punchline у червоній рамці знизу
    _shape_rounded(slide, Cm(2), Cm(10.3), Cm(21.4), Cm(2.5),
                   fill=COLOR_ACCENT_2)
    _text(slide, Cm(2), Cm(10.3), Cm(21.4), Cm(2.5),
          data["punchline"], size=18, bold=True,
          color=COLOR_WHITE, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 4 – split problem

def _slide_problem_split(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    _text(slide, Cm(1), Cm(1.0), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)
    _text(slide, Cm(1), Cm(1.8), prs.slide_width - Cm(2), Cm(1.2),
          data["title"], size=28, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

    # дві колонки-картки
    col_w = Cm(11.5)
    col_h = Cm(7.5)
    col_y = Cm(3.8)
    # ліва
    _shape_rounded(slide, Cm(1), col_y, col_w, col_h, fill=COLOR_WHITE)
    _shape_rect(slide, Cm(1), col_y, Cm(0.3), col_h, fill=COLOR_ACCENT_1)
    _text(slide, Cm(1.6), col_y + Cm(0.5), col_w - Cm(1), Cm(1.2),
          data["left_title"], size=18, bold=True, color=COLOR_ACCENT_1)
    _text(slide, Cm(1.6), col_y + Cm(2.0), col_w - Cm(1), col_h - Cm(2.5),
          data["left_text"], size=14, color=COLOR_TEXT,
          align=PP_ALIGN.LEFT)

    # права
    _shape_rounded(slide, Cm(13), col_y, col_w, col_h, fill=COLOR_WHITE)
    _shape_rect(slide, Cm(13), col_y, Cm(0.3), col_h, fill=COLOR_ACCENT_2)
    _text(slide, Cm(13.6), col_y + Cm(0.5), col_w - Cm(1), Cm(1.2),
          data["right_title"], size=18, bold=True, color=COLOR_ACCENT_2)
    _text(slide, Cm(13.6), col_y + Cm(2.0), col_w - Cm(1), col_h - Cm(2.5),
          data["right_text"], size=14, color=COLOR_TEXT,
          align=PP_ALIGN.LEFT)

    # footer цитата
    _shape_rect(slide, Cm(1), Cm(12), Cm(23.4), Cm(0.05),
                fill=COLOR_PRIMARY)
    _text(slide, Cm(1), Cm(12.2), Cm(23.4), Cm(1.2),
          data["footer"], size=18, bold=True, italic=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 5 i 6 – two col card

def _slide_two_col_card(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    _text(slide, Cm(1), Cm(1.0), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_1, align=PP_ALIGN.LEFT)
    _text(slide, Cm(1), Cm(1.8), prs.slide_width - Cm(2), Cm(1.5),
          data["title"], size=32, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

    # ліва картка з буллетами
    _shape_rounded(slide, Cm(1), Cm(4.2), Cm(13), Cm(8.5),
                   fill=COLOR_WHITE)
    _shape_rect(slide, Cm(1), Cm(4.2), Cm(0.3), Cm(8.5),
                fill=COLOR_ACCENT_1)
    _bullet_list(slide, Cm(1.7), Cm(5), Cm(12), Cm(7),
                 data["left_bullets"], size=18,
                 marker_color=COLOR_ACCENT_1, marker="●")

    # права картка з цитатою на акцентному фоні
    _shape_rounded(slide, Cm(14.6), Cm(4.2), Cm(9.8), Cm(8.5),
                   fill=COLOR_PRIMARY)
    _text(slide, Cm(14.6), Cm(4.2), Cm(9.8), Cm(8.5),
          f"“{data['right_quote']}”", size=22, bold=True,
          color=COLOR_WHITE, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 7 – INSIGHT dark

def _slide_insight_dark(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # повністю темний фон
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_DARK)
    # маленький жовтий акцент зверху
    _shape_rect(slide, Cm(0), Cm(0), Cm(1.5), Cm(0.4),
                fill=COLOR_ACCENT_3)
    _text(slide, Cm(1), Cm(1), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_3, align=PP_ALIGN.LEFT)

    # головна цитата по центру
    _text(slide, Cm(1.5), Cm(3), Cm(22.4), Cm(4),
          f"“{data['big_quote']}”", size=28, bold=True,
          color=COLOR_WHITE, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)

    # роздільна лінія
    _shape_rect(slide, Cm(11), Cm(7.5), Cm(3.4), Cm(0.08),
                fill=COLOR_ACCENT_2)

    # підцитата
    _text(slide, Cm(1.5), Cm(8), Cm(22.4), Cm(2.5),
          data["subtitle_quote"], size=16, italic=True,
          color=RGBColor(0xBB, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    # назва моделі великими літерами знизу
    _shape_rounded(slide, Cm(9), Cm(11.5), Cm(7.4), Cm(1.5),
                   fill=COLOR_ACCENT_2)
    _text(slide, Cm(9), Cm(11.5), Cm(7.4), Cm(1.5),
          data["model_name"], size=24, bold=True,
          color=COLOR_WHITE, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE, font="Consolas")

    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 8 – architecture

def _slide_architecture(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    _text(slide, Cm(1), Cm(1.0), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)
    _text(slide, Cm(1), Cm(1.8), prs.slide_width - Cm(2), Cm(1.2),
          data["title"], size=28, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

    # ліворуч велика картинка з тінню
    if Path(data["image"]).exists():
        _shadow_image(slide, data["image"],
                      Cm(1), Cm(3.5), Cm(15), Cm(9))

    # праворуч – статистика-картки
    stat_x = Cm(17)
    stat_w = Cm(7.4)
    n = len(data["right_stats"])
    card_h = Cm(2.0)
    gap = Cm(0.3)
    start_y = Cm(3.5)
    for i, (big, label) in enumerate(data["right_stats"]):
        y = start_y + (card_h + gap) * i
        _shape_rounded(slide, stat_x, y, stat_w, card_h, fill=COLOR_WHITE)
        _shape_rect(slide, stat_x, y, Cm(0.25), card_h, fill=COLOR_ACCENT_2)
        _text(slide, stat_x + Cm(0.5), y, Cm(3), card_h,
              big, size=28, bold=True, color=COLOR_ACCENT_2,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, stat_x + Cm(3.5), y, stat_w - Cm(3.5), card_h,
              label, size=13, color=COLOR_TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 9 – method pipeline

def _slide_method_pipeline(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    _text(slide, Cm(1), Cm(1.0), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)
    _text(slide, Cm(1), Cm(1.8), prs.slide_width - Cm(2), Cm(1.2),
          data["title"], size=28, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

    # 4 кроки pipeline у ряд
    steps = data["pipeline_steps"]
    n = len(steps)
    step_w = Cm(5.4)
    gap = Cm(0.6)
    total_w = step_w * n + gap * (n - 1)
    start_x = (prs.slide_width - total_w) // 2
    step_y = Cm(4.5)

    for i, (num, name, desc) in enumerate(steps):
        x = start_x + (step_w + gap) * i
        # тінь
        _shape_rounded(slide, x + Cm(0.15), step_y + Cm(0.15),
                       step_w, Cm(6.5), fill=COLOR_SHADOW)
        # картка
        _shape_rounded(slide, x, step_y, step_w, Cm(6.5),
                       fill=COLOR_WHITE)
        # темна верхня смуга з номером
        _shape_rect(slide, x, step_y, step_w, Cm(1.3),
                    fill=COLOR_PRIMARY)
        _text(slide, x, step_y, step_w, Cm(1.3),
              num, size=22, bold=True, color=COLOR_ACCENT_3,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
              font="Consolas")
        # назва кроку
        _text(slide, x + Cm(0.3), step_y + Cm(1.8), step_w - Cm(0.6),
              Cm(1.5),
              name, size=16, bold=True, color=COLOR_PRIMARY,
              align=PP_ALIGN.CENTER)
        # опис
        _text(slide, x + Cm(0.3), step_y + Cm(3.5), step_w - Cm(0.6),
              Cm(2),
              desc, size=12, color=COLOR_MUTED,
              align=PP_ALIGN.CENTER)
        # стрілка між картками
        if i < n - 1:
            arrow_x = x + step_w - Cm(0.2)
            _text(slide, arrow_x, step_y + Cm(2.5), Cm(1), Cm(2),
                  "❯", size=24, bold=True, color=COLOR_ACCENT_2,
                  align=PP_ALIGN.CENTER)

    # tagline знизу
    _text(slide, Cm(1), Cm(12.2), prs.slide_width - Cm(2), Cm(0.8),
          data["tagline"], size=14, italic=True,
          color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 10 – podium

def _slide_podium(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    _text(slide, Cm(1), Cm(1.0), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)
    _text(slide, Cm(1), Cm(1.8), prs.slide_width - Cm(2), Cm(1.2),
          data["title"], size=28, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

    # три картки на подіумі
    podium = data["podium"]
    card_w = Cm(7.5)
    gap = Cm(0.4)
    total_w = card_w * 3 + gap * 2
    start_x = (prs.slide_width - total_w) // 2
    card_y = Cm(4)
    card_h = Cm(8)

    for i, p in enumerate(podium):
        x = start_x + (card_w + gap) * i
        # тінь
        _shape_rounded(slide, x + Cm(0.15), card_y + Cm(0.15),
                       card_w, card_h, fill=COLOR_SHADOW)
        # картка
        _shape_rounded(slide, x, card_y, card_w, card_h, fill=COLOR_WHITE)
        # кольорова верхня смуга
        _shape_rect(slide, x, card_y, card_w, Cm(1.3), fill=p["color"])
        # назва датасета
        _text(slide, x, card_y, card_w, Cm(1.3),
              p["rank"], size=16, bold=True, color=COLOR_WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # ім'я моделі-переможця крупно
        _text(slide, x + Cm(0.3), card_y + Cm(1.8), card_w - Cm(0.6),
              Cm(2.3),
              p["winner"], size=26, bold=True, color=COLOR_PRIMARY,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
              font="Consolas")
        # score
        _shape_rounded(slide, x + Cm(1.5), card_y + Cm(4.5),
                       card_w - Cm(3), Cm(1),
                       fill=COLOR_ACCENT_3)
        _text(slide, x + Cm(1.5), card_y + Cm(4.5),
              card_w - Cm(3), Cm(1),
              p["score"], size=15, bold=True, color=COLOR_PRIMARY,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # коментар
        _text(slide, x + Cm(0.3), card_y + Cm(5.8), card_w - Cm(0.6),
              Cm(2),
              p["comment"], size=11, color=COLOR_MUTED,
              align=PP_ALIGN.CENTER)

    # ключовий висновок знизу
    _shape_rect(slide, Cm(1), Cm(12.5), Cm(23.4), Cm(0.05),
                fill=COLOR_ACCENT_2)
    _text(slide, Cm(1), Cm(12.6), Cm(23.4), Cm(0.8),
          data["insight"], size=14, bold=True, italic=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 11 – contribution

def _slide_contribution(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    _text(slide, Cm(1), Cm(1.0), Cm(15), Cm(0.8),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)
    _text(slide, Cm(1), Cm(1.8), prs.slide_width - Cm(2), Cm(1.5),
          data["title"], size=28, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

    # story block – велика цитата у рамці
    _shape_rounded(slide, Cm(1), Cm(4), Cm(23.4), Cm(2.8),
                   fill=COLOR_WHITE, line_color=COLOR_ACCENT_2)
    _text(slide, Cm(1.5), Cm(4), Cm(22.4), Cm(2.8),
          data["story_block"], size=15, italic=True,
          color=COLOR_TEXT, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)

    # стрілка/розділювач вниз
    _text(slide, Cm(11.5), Cm(7), Cm(2.4), Cm(1),
          "↓", size=28, bold=True, color=COLOR_ACCENT_1,
          align=PP_ALIGN.CENTER)

    # fix block – темна картка
    _shape_rounded(slide, Cm(1), Cm(8.2), Cm(23.4), Cm(3.5),
                   fill=COLOR_PRIMARY)
    _text(slide, Cm(1.5), Cm(8.4), Cm(7), Cm(0.8),
          data["fix_label"], size=14, bold=True,
          color=COLOR_ACCENT_3, align=PP_ALIGN.LEFT)
    _text(slide, Cm(1.5), Cm(9.2), Cm(22.4), Cm(2.5),
          data["fix_text"], size=14,
          color=COLOR_WHITE, align=PP_ALIGN.LEFT)

    # result
    _text(slide, Cm(1), Cm(12.1), Cm(23.4), Cm(0.9),
          f"✓  {data['result_text']}", size=16, bold=True,
          color=COLOR_ACCENT_1, align=PP_ALIGN.CENTER)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


# =============================================================================
# SLIDE 12 – final QR

def _slide_image_with_caption(prs, data):
    """Слайд із зображенням (зокрема скріншотом таблиці) на повну ширину.

    Очікувані поля: label, title, image (Path), tagline (опційно).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_BG_LIGHT)
    _act_badge(slide, prs, data["act"])

    # маленький лейбл вгорі
    _text(slide, Cm(1), Cm(0.7), Cm(20), Cm(0.7),
          data["label"].upper(), size=12, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT)

    # title
    _text(slide, Cm(1), Cm(1.5), prs.slide_width - Cm(2), Cm(1.2),
          data["title"], size=24, bold=True,
          color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

    # картинка по центру
    img = data.get("image")
    if img and Path(img).exists():
        # рамка-тінь
        max_w = Cm(20)
        max_h = Cm(9.5)
        x = (prs.slide_width - max_w) // 2
        y = Cm(3.0)
        _shape_rect(slide, x + Cm(0.15), y + Cm(0.15),
                    max_w, max_h, fill=COLOR_SHADOW)
        _shape_rect(slide, x, y, max_w, max_h, fill=COLOR_WHITE)
        slide.shapes.add_picture(str(img),
                                 x + Cm(0.2), y + Cm(0.2),
                                 width=max_w - Cm(0.4),
                                 height=max_h - Cm(0.4))

    # підпис знизу
    tagline = data.get("tagline")
    if tagline:
        _text(slide, Cm(1), Cm(13), prs.slide_width - Cm(2), Cm(0.8),
              tagline, size=13, italic=True,
              color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    _footer_brand(slide, prs)
    _add_notes(slide, data["speech"])


def _slide_thanks_big(prs, data):
    """Великий слайд «Дякую за увагу!» – головне посередині на всю ширину."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # темний фон
    _shape_rect(slide, Cm(0), Cm(0), prs.slide_width, prs.slide_height,
                fill=COLOR_PRIMARY)
    # 3 акцентні смуги по верху і низу
    for j, col in enumerate(
            (COLOR_ACCENT_2, COLOR_ACCENT_1, COLOR_ACCENT_3)):
        _shape_rect(slide, Cm(0.5 + j * 1.6), Cm(0.5),
                    Cm(1.3), Cm(0.25), fill=col)
        _shape_rect(slide,
                    prs.slide_width - Cm(1.8 + j * 1.6),
                    prs.slide_height - Cm(0.75),
                    Cm(1.3), Cm(0.25), fill=col)

    # Дуже велике "Дякую за увагу!" посередині
    _text(slide, Cm(0), Cm(3.5), prs.slide_width, Cm(4),
          data["title"], size=96, bold=True,
          color=COLOR_WHITE, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE)

    # підзаголовок під ним
    _text(slide, Cm(0), Cm(7.5), prs.slide_width, Cm(1.2),
          data["subtitle"], size=22, italic=True,
          color=COLOR_ACCENT_3, align=PP_ALIGN.CENTER)

    # GitHub label + URL
    if data.get("github_label"):
        _text(slide, Cm(0), Cm(9.6), prs.slide_width, Cm(0.7),
              data["github_label"], size=14,
              color=RGBColor(0xBB, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    if data.get("github_url"):
        _text(slide, Cm(0), Cm(10.3), prs.slide_width, Cm(0.9),
              data["github_url"], size=18, bold=True,
              color=COLOR_ACCENT_2, align=PP_ALIGN.CENTER,
              font="Consolas")

    # QR справа знизу (маленький)
    qr = data.get("qr_image")
    if qr and Path(qr).exists():
        qr_size = Cm(3.0)
        qr_x = prs.slide_width - qr_size - Cm(0.8)
        qr_y = prs.slide_height - qr_size - Cm(1.2)
        _shape_rect(slide, qr_x, qr_y, qr_size, qr_size, fill=COLOR_WHITE)
        slide.shapes.add_picture(str(qr),
                                 qr_x + Cm(0.2), qr_y + Cm(0.2),
                                 width=qr_size - Cm(0.4),
                                 height=qr_size - Cm(0.4))

    _add_notes(slide, data["speech"])


def _slide_final_qr(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # ліва половина темна
    _shape_rect(slide, Cm(0), Cm(0), Cm(13), prs.slide_height,
                fill=COLOR_PRIMARY)
    # права світла
    _shape_rect(slide, Cm(13), Cm(0), prs.slide_width - Cm(13),
                prs.slide_height, fill=COLOR_BG_LIGHT)

    # 3 акцентні смуги
    _shape_rect(slide, Cm(0), Cm(0.5), Cm(2.0), Cm(0.2),
                fill=COLOR_ACCENT_2)
    _shape_rect(slide, Cm(0), Cm(0.9), Cm(1.3), Cm(0.2),
                fill=COLOR_ACCENT_3)
    _shape_rect(slide, Cm(0), Cm(1.3), Cm(0.7), Cm(0.2),
                fill=COLOR_ACCENT_1)

    # лейбл
    _text(slide, Cm(1), Cm(2), Cm(11), Cm(0.8),
          data["label"].upper(), size=13, bold=True,
          color=COLOR_ACCENT_3, align=PP_ALIGN.LEFT)

    # title
    _text(slide, Cm(1), Cm(3.5), Cm(11), Cm(2),
          data["title"], size=32, bold=True,
          color=COLOR_WHITE, align=PP_ALIGN.LEFT)

    # github label
    _text(slide, Cm(1), Cm(7), Cm(11), Cm(0.8),
          data["github_label"], size=14,
          color=RGBColor(0xBB, 0xCC, 0xCC), align=PP_ALIGN.LEFT)

    # github URL
    _text(slide, Cm(1), Cm(7.8), Cm(11), Cm(1),
          data["github_url"], size=14, bold=True,
          color=COLOR_ACCENT_3, align=PP_ALIGN.LEFT,
          font="Consolas")

    # велике дякую внизу
    _text(slide, Cm(1), Cm(11), Cm(11), Cm(2),
          data["thanks"], size=36, bold=True,
          color=COLOR_ACCENT_2, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.BOTTOM)

    # QR code справа
    qr_x = Cm(15)
    qr_y = Cm(2.5)
    qr_size = Cm(8.5)
    qr = data.get("qr_image")
    if qr and Path(qr).exists():
        _shape_rect(slide, qr_x + Cm(0.2), qr_y + Cm(0.2),
                    qr_size, qr_size, fill=COLOR_SHADOW)
        _shape_rect(slide, qr_x, qr_y, qr_size, qr_size,
                    fill=COLOR_WHITE)
        slide.shapes.add_picture(str(qr),
                                 qr_x + Cm(0.3), qr_y + Cm(0.3),
                                 width=qr_size - Cm(0.6),
                                 height=qr_size - Cm(0.6))
    # caption під QR
    _text(slide, qr_x, qr_y + qr_size + Cm(0.3),
          qr_size, Cm(1),
          data["qr_caption"], size=12, italic=True,
          color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    _add_notes(slide, data["speech"])


# =============================================================================

def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.29)

    builders = {
        "title_hero": _slide_title_hero,
        "big_numbers": _slide_big_numbers,
        "problem_question": _slide_problem_question,
        "problem_split": _slide_problem_split,
        "two_col_card": _slide_two_col_card,
        "insight_dark": _slide_insight_dark,
        "architecture_slide": _slide_architecture,
        "method_pipeline": _slide_method_pipeline,
        "podium_results": _slide_podium,
        "contribution": _slide_contribution,
        "final_qr": _slide_final_qr,
        "thanks_big": _slide_thanks_big,
        "image_with_caption": _slide_image_with_caption,
    }
    for data in SLIDES:
        builders[data["kind"]](prs, data)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_PATH))
    return PPTX_PATH


# =============================================================================
# DOCX – speech

def _docx_set_run(run, *, size=12, bold=False, italic=False,
                  font="Times New Roman"):
    run.font.name = font
    run.font.size = DPt(size)
    run.font.bold = bold
    run.font.italic = italic
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    rfonts.set(qn("w:ascii"), font)
    rfonts.set(qn("w:hAnsi"), font)
    rfonts.set(qn("w:cs"), font)


def _docx_para(doc, text, *, size=12, bold=False, italic=False,
               align=None, indent=DCm(1.25), space=DPt(6)):
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    pf = p.paragraph_format
    if indent is not None:
        pf.first_line_indent = indent
    pf.space_after = space
    pf.line_spacing = 1.5
    run = p.add_run(text)
    _docx_set_run(run, size=size, bold=bold, italic=italic)


def _docx_heading(doc, text, level=1):
    p = doc.add_paragraph()
    pf = p.paragraph_format
    pf.space_before = DPt(14)
    pf.space_after = DPt(8)
    pf.keep_with_next = True
    pf.line_spacing = 1.15
    p.alignment = DOCX_ALIGN.CENTER if level == 1 else DOCX_ALIGN.LEFT
    run = p.add_run(text)
    _docx_set_run(run, size=16 if level == 1 else 13, bold=True)


def build_docx() -> Path:
    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = DPt(12)
    for section in doc.sections:
        section.left_margin = DCm(2.5)
        section.right_margin = DCm(1.5)
        section.top_margin = DCm(2)
        section.bottom_margin = DCm(2)

    for _ in range(4):
        doc.add_paragraph()
    _docx_para(doc, "ТЕКСТ ДОПОВІДІ", size=20, bold=True,
               align=DOCX_ALIGN.CENTER, indent=None, space=DPt(8))
    _docx_para(doc, "Сторітелінгова презентація (12 слайдів, 8-10 хв)",
               size=14, italic=True, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(20))
    _docx_para(doc, "12 моделей. 3 датасети. 1 ідея.",
               size=16, bold=True, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(60))
    _docx_para(doc, "Волошенюк Богдан Анатолійович, група 441-А",
               size=12, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(2))
    _docx_para(doc, "ЧНУ ім. Юрія Федьковича, 2026",
               size=12, align=DOCX_ALIGN.CENTER,
               indent=None, space=DPt(20))
    doc.add_page_break()

    _docx_heading(doc, "Як влаштовано виступ", level=1)
    _docx_para(doc,
        "Презентація побудована за класичною сторітелінговою структурою "
        "з п'ятьох актів: HOOK – щоб зачепити, PROBLEM – щоб "
        "сформулювати задачу, INSIGHT – щоб подати ключову ідею, "
        "SOLUTION – щоб показати рішення, CONCLUSION – щоб залишити "
        "сильне останнє враження.")
    _docx_para(doc,
        "Тривалість – приблизно вісім-десять хвилин при спокійному темпі. "
        "На кожен слайд відведено від тридцяти секунд до хвилини. На "
        "слайдах з аха-моментом (слайд 7) і подіумом (слайд 10) "
        "рекомендується робити паузу – щоб комісія встигла "
        "сприйняти ключову думку.")
    _docx_para(doc,
        "На захисті використовувати Режим доповідача в PowerPoint – "
        "нотатки видно лише вам, на основному екрані – сам слайд. На "
        "останньому слайді – QR-код. Можна запропонувати членам "
        "комісії сканувати його зі смартфона.")
    doc.add_page_break()

    _docx_heading(doc, "Структура виступу", level=1)
    acts = [
        ("HOOK", "Слайди 1-2", "Зачепити увагу масштабом."),
        ("PROBLEM", "Слайди 3-4", "Описати проблему ML і гіперпараметри."),
        ("INSIGHT", "Слайди 5-7", "Розкрити теорію і ключову ідею."),
        ("SOLUTION", "Слайди 8-10", "Показати реалізацію і результати."),
        ("CONCLUSION", "Слайди 11-12", "Власний внесок і подяка."),
    ]
    for act, slides_range, descr in acts:
        p = doc.add_paragraph()
        pf = p.paragraph_format
        pf.left_indent = DCm(0.75)
        pf.space_after = DPt(4)
        r1 = p.add_run(f"{act} ({slides_range}) – ")
        _docx_set_run(r1, size=12, bold=True)
        r2 = p.add_run(descr)
        _docx_set_run(r2, size=12)
    doc.add_page_break()

    _docx_heading(doc, "Суцільний текст доповіді", level=1)
    _docx_para(doc,
        "Версія для тих, хто хоче читати все підряд без перемикання "
        "між слайдами. Між абзацами доцільно робити паузу 2-3 секунди.",
        italic=True)
    for data in SLIDES:
        _docx_para(doc, data["speech"], size=12)
    doc.add_page_break()

    _docx_heading(doc, "Текст по слайдах", level=1)
    for i, data in enumerate(SLIDES, start=1):
        title = data.get("title", data.get("label", "")).replace("\n", " ")
        _docx_heading(doc, f"Слайд {i} ({data['act']}). {title}", level=2)
        _docx_para(doc, "Що читати:", bold=True,
                   indent=None, space=DPt(2))
        _docx_para(doc, data["speech"], size=12)
        doc.add_paragraph()

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(DOCX_PATH))
    return DOCX_PATH


def main() -> int:
    print(f"Генерую PowerPoint ({len(SLIDES)} слайдів)...")
    pptx_out = build_pptx()
    print(f"  OK: {pptx_out} ({pptx_out.stat().st_size // 1024} KB)")
    print("Генерую текст доповіді (.docx)...")
    docx_out = build_docx()
    print(f"  OK: {docx_out} ({docx_out.stat().st_size // 1024} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
